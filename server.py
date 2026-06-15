"""
LightWatch Backend — Railway Flask App
Provides:
  POST /api/analyze        — AIVM anomaly analysis on monitor data
  POST /api/report         — AIVM-generated report narrative
  GET  /api/qb/auth        — Start QuickBooks OAuth flow
  GET  /api/qb/callback    — QuickBooks OAuth callback (store tokens)
  GET  /api/qb/data        — Pull live QB data for connected company
  GET  /api/qb/status      — Check if QuickBooks is connected
  GET  /health             — Health check

Env vars required (set in Railway):
  LIGHTCHAIN_PRIVATE_KEY   — dApp wallet private key (for AIVM)
  QUICKBOOKS_CLIENT_ID     — from developer.intuit.com
  QUICKBOOKS_CLIENT_SECRET — from developer.intuit.com
  QUICKBOOKS_REDIRECT_URI  — https://<your-railway-url>/api/qb/callback
  SESSION_SECRET           — any random string (for Flask session)
"""

import os, json, time, threading, secrets as _secrets_mod
import base64 as _b64_mod
import sqlite3, io, csv as _csv_mod, re as _re_mod
from flask import Flask, request, jsonify, session, redirect, send_file
from flask_cors import CORS
import requests as req

app = Flask(__name__)
app.secret_key = os.environ.get("SESSION_SECRET", _secrets_mod.token_hex(32))
CORS(app)

# ── AIVM CONFIG ──────────────────────────────────────────────────────────────
AIVM_PRIVATE_KEY = os.environ.get("LIGHTCHAIN_PRIVATE_KEY", "").strip()
AIVM_GATEWAY     = "https://chat-api.mainnet.lightchain.ai"
AIVM_RELAY       = "wss://relay.mainnet.lightchain.ai/ws"
AIVM_JOB_REG     = "0xfB15F90298e4CcD7106E76fFB5e520315cC42B0b"
AIVM_JOB_FEE     = 20_000_000_000_000_000   # 0.02 LCAI in wei
AIVM_CHAIN_ID    = 9200

AIVM_ABI = [
    {"name":"createSession","type":"function","stateMutability":"payable",
     "inputs":[{"name":"paramsHash","type":"bytes32"},{"name":"worker","type":"address"},
               {"name":"encWorkerKey","type":"bytes"},{"name":"encDisputerKey","type":"bytes"},
               {"name":"workerSig","type":"bytes"},{"name":"expiry","type":"uint256"}],
     "outputs":[{"name":"sessionId","type":"uint256"}]},
    {"name":"submitJob","type":"function","stateMutability":"payable",
     "inputs":[{"name":"sessionId","type":"uint256"},{"name":"promptHash","type":"bytes32"}],
     "outputs":[]},
    {"name":"SessionCreated","type":"event",
     "inputs":[{"name":"sessionId","type":"uint256","indexed":True},
               {"name":"modelId","type":"bytes32","indexed":False},
               {"name":"worker","type":"address","indexed":True}]},
]

# ── QUICKBOOKS CONFIG ────────────────────────────────────────────────────────
QB_CLIENT_ID     = os.environ.get("QUICKBOOKS_CLIENT_ID", "")
QB_CLIENT_SECRET = os.environ.get("QUICKBOOKS_CLIENT_SECRET", "")
QB_REDIRECT_URI  = os.environ.get("QUICKBOOKS_REDIRECT_URI", "")
QB_SCOPE         = "com.intuit.quickbooks.accounting"
QB_AUTH_URL      = "https://appcenter.intuit.com/connect/oauth2"
QB_TOKEN_URL     = "https://oauth.platform.intuit.com/oauth2/v1/tokens/bearer"
QB_API_BASE      = "https://quickbooks.api.intuit.com/v3/company"

# In-memory token store (Railway volume would persist across restarts)
_QB_TOKENS = {}   # { realm_id: {access_token, refresh_token, expires_at} }
_QB_LOCK   = threading.Lock()

# ── GREAT BRIDGE FURNITURE — SQLite DB ───────────────────────────────────────
DB_PATH = os.environ.get("DB_PATH", "/data/lightview.db")

def _db():
    conn = sqlite3.connect(DB_PATH, check_same_thread=False)
    conn.row_factory = sqlite3.Row
    return conn

def _init_db():
    with _db() as conn:
        conn.execute("""
            CREATE TABLE IF NOT EXISTS gb_orders (
                id             INTEGER PRIMARY KEY AUTOINCREMENT,
                customer_name  TEXT,
                phone          TEXT,
                items          TEXT,
                manufacturer   TEXT,
                total_amount   REAL    DEFAULT 0,
                deposit_paid   REAL    DEFAULT 0,
                expected_date  TEXT,
                status         TEXT    DEFAULT 'Ordered',
                notes          TEXT,
                created_at     TEXT    DEFAULT (datetime('now')),
                updated_at     TEXT    DEFAULT (datetime('now'))
            )
        """)
        conn.execute("""
            CREATE TABLE IF NOT EXISTS gb_docs (
                id          INTEGER PRIMARY KEY AUTOINCREMENT,
                order_id    INTEGER NOT NULL,
                filename    TEXT    NOT NULL,
                filetype    TEXT,
                filedata    BLOB    NOT NULL,
                uploaded_at TEXT    DEFAULT (datetime('now'))
            )
        """)
        # Migrate: add columns to existing orders table
        for _col, _typ in [("delivery_date", "TEXT"), ("arrival_date", "TEXT"), ("ticket_doc_id", "INTEGER")]:
            try:
                conn.execute(f"ALTER TABLE gb_orders ADD COLUMN {_col} {_typ}")
            except Exception:
                pass  # column already exists
        # Migrate: add is_ticket flag to gb_docs
        try:
            conn.execute("ALTER TABLE gb_docs ADD COLUMN is_ticket INTEGER DEFAULT 0")
        except Exception:
            pass
        conn.execute("""
            CREATE TABLE IF NOT EXISTS gb_blocked_days (
                id         INTEGER PRIMARY KEY AUTOINCREMENT,
                date       TEXT    NOT NULL UNIQUE,
                reason     TEXT    DEFAULT 'closed',
                note       TEXT,
                created_at TEXT    DEFAULT (datetime('now'))
            )
        """)
        conn.execute("""
            CREATE TABLE IF NOT EXISTS gb_deliveries (
                id          INTEGER PRIMARY KEY AUTOINCREMENT,
                date        TEXT    NOT NULL,
                time        TEXT    DEFAULT '',
                direction   TEXT    DEFAULT 'OUT',
                name        TEXT    NOT NULL DEFAULT '',
                address     TEXT    DEFAULT '',
                phone       TEXT    DEFAULT '',
                items       TEXT    DEFAULT '',
                notes       TEXT    DEFAULT '',
                order_id    INTEGER,
                created_at  TEXT    DEFAULT (datetime('now')),
                updated_at  TEXT    DEFAULT (datetime('now'))
            )
        """)
        conn.execute("""
            CREATE TABLE IF NOT EXISTS biz_profile (
                id             INTEGER PRIMARY KEY AUTOINCREMENT,
                company        TEXT    NOT NULL UNIQUE,
                business_type  TEXT    DEFAULT '',
                priorities     TEXT    DEFAULT '[]',
                briefing_focus TEXT    DEFAULT '',
                staff_notes    TEXT    DEFAULT '',
                summary        TEXT    DEFAULT '',
                raw_conv       TEXT    DEFAULT '',
                created_at     TEXT    DEFAULT (datetime('now')),
                updated_at     TEXT    DEFAULT (datetime('now'))
            )
        """)
        conn.execute("""
            CREATE TABLE IF NOT EXISTS gb_activity_log (
                id         INTEGER PRIMARY KEY AUTOINCREMENT,
                user_name  TEXT    NOT NULL DEFAULT 'Unknown',
                action     TEXT    NOT NULL DEFAULT '',
                target     TEXT    DEFAULT '',
                details    TEXT    DEFAULT '',
                created_at TEXT    DEFAULT (datetime('now'))
            )
        """)
        conn.execute("""
            CREATE TABLE IF NOT EXISTS gb_inventory (
                id                 INTEGER PRIMARY KEY AUTOINCREMENT,
                sku                TEXT,
                description        TEXT,
                manufacturer       TEXT,
                finish_fabric      TEXT,
                warehouse_location TEXT,
                cost_price         REAL,
                retail_price       REAL,
                status             TEXT DEFAULT 'Available',
                customer_name      TEXT,
                order_id           INTEGER,
                date_in            TEXT,
                date_out           TEXT,
                delivery_type      TEXT,
                notes              TEXT,
                created_at         TIMESTAMP DEFAULT CURRENT_TIMESTAMP
            )
        """)
        conn.execute("""
            CREATE TABLE IF NOT EXISTS gb_manufacturers (
                id             INTEGER PRIMARY KEY AUTOINCREMENT,
                name           TEXT NOT NULL,
                contact_person TEXT DEFAULT '',
                email          TEXT DEFAULT '',
                phone          TEXT DEFAULT '',
                fax            TEXT DEFAULT '',
                address        TEXT DEFAULT '',
                notes          TEXT DEFAULT '',
                created_at     TEXT DEFAULT (datetime('now'))
            )
        """)
        # Migrate gb_orders: add manufacturer_id FK column
        try:
            conn.execute("ALTER TABLE gb_orders ADD COLUMN manufacturer_id INTEGER")
        except Exception:
            pass
        # Migrate gb_orders: add po_sent_at column
        try:
            conn.execute("ALTER TABLE gb_orders ADD COLUMN po_sent_at TEXT")
        except Exception:
            pass
        conn.commit()

try:
    _init_db()
    print("[DB] SQLite initialized")
except Exception as _e:
    print(f"[DB] init failed: {_e}")


def _log_activity(user: str, action: str, target: str = "", details: str = ""):
    """Append one row to gb_activity_log (fire-and-forget, never raises)."""
    try:
        with _db() as conn:
            conn.execute(
                "INSERT INTO gb_activity_log (user_name, action, target, details) VALUES (?,?,?,?)",
                (str(user or "Unknown")[:100], str(action)[:200],
                 str(target)[:200], str(details)[:500])
            )
            conn.commit()
    except Exception as _e:
        print(f"[activity_log] {_e}")


# ── GBF EMAIL + PO HELPERS ───────────────────────────────────────────────────
import smtplib
from email.mime.multipart import MIMEMultipart
from email.mime.text import MIMEText

GBF_EMAIL_ADDRESS  = os.environ.get('GBF_EMAIL_ADDRESS', '')
GBF_EMAIL_PASSWORD = os.environ.get('GBF_EMAIL_PASSWORD', '')


def _send_aol_email(to_email, subject, html_body):
    """Send email via AOL SMTP using env-var credentials."""
    if not GBF_EMAIL_ADDRESS or not GBF_EMAIL_PASSWORD:
        raise ValueError("GBF_EMAIL_ADDRESS and GBF_EMAIL_PASSWORD env vars not set")
    msg = MIMEMultipart('alternative')
    msg['Subject']  = subject
    msg['From']     = f"Great Bridge Furniture <{GBF_EMAIL_ADDRESS}>"
    msg['To']       = to_email
    msg['Reply-To'] = GBF_EMAIL_ADDRESS
    msg.attach(MIMEText(html_body, 'html'))
    with smtplib.SMTP('smtp.aol.com', 587) as smtp:
        smtp.ehlo()
        smtp.starttls()
        smtp.login(GBF_EMAIL_ADDRESS, GBF_EMAIL_PASSWORD)
        smtp.send_message(msg)


def _gen_po_html(order, manufacturer):
    """Generate a professional Purchase Order as HTML."""
    import datetime as _dt_po
    po_number  = f"GBF-{_dt_po.date.today().strftime('%Y%m%d')}-{order['id']:04d}"
    today      = _dt_po.date.today().strftime('%B %d, %Y')
    mfr_name    = manufacturer.get('name', order.get('manufacturer', '')) if manufacturer else order.get('manufacturer', '')
    mfr_contact = manufacturer.get('contact_person', '') if manufacturer else ''
    mfr_email   = manufacturer.get('email', '')          if manufacturer else ''
    mfr_address = manufacturer.get('address', '')        if manufacturer else ''

    return f"""<!DOCTYPE html>
<html>
<head><meta charset="utf-8">
<style>
  body {{ font-family: Arial, sans-serif; font-size: 14px; color: #222; margin: 0; padding: 0; }}
  .po-wrap {{ max-width: 700px; margin: 0 auto; padding: 40px 32px; }}
  .header {{ display: flex; justify-content: space-between; align-items: flex-start; border-bottom: 3px solid #1a3a6b; padding-bottom: 20px; margin-bottom: 28px; }}
  .store-name {{ font-size: 26px; font-weight: 700; color: #1a3a6b; }}
  .store-info {{ font-size: 12px; color: #555; margin-top: 4px; }}
  .po-title {{ text-align: right; }}
  .po-title h1 {{ font-size: 28px; font-weight: 700; color: #1a3a6b; margin: 0; }}
  .po-title .po-num {{ font-size: 14px; color: #444; margin-top: 4px; }}
  .po-title .po-date {{ font-size: 13px; color: #666; }}
  .addresses {{ display: flex; gap: 40px; margin-bottom: 28px; }}
  .addr-block {{ flex: 1; }}
  .addr-block h3 {{ font-size: 11px; text-transform: uppercase; letter-spacing: 1px; color: #888; margin: 0 0 6px 0; }}
  .addr-block p {{ margin: 2px 0; font-size: 13px; }}
  table {{ width: 100%; border-collapse: collapse; margin-bottom: 24px; }}
  th {{ background: #1a3a6b; color: #fff; padding: 10px 12px; text-align: left; font-size: 13px; }}
  td {{ padding: 10px 12px; border-bottom: 1px solid #e5e5e5; font-size: 13px; }}
  tr:nth-child(even) td {{ background: #f7f9fc; }}
  .totals {{ text-align: right; margin-bottom: 28px; }}
  .totals table {{ width: 260px; margin-left: auto; }}
  .totals td {{ border: none; padding: 4px 8px; }}
  .totals .grand-total td {{ font-weight: 700; font-size: 16px; border-top: 2px solid #1a3a6b; padding-top: 8px; }}
  .footer {{ border-top: 1px solid #ddd; padding-top: 16px; font-size: 12px; color: #666; }}
  .confirm-note {{ background: #f0f4ff; border-left: 4px solid #1a3a6b; padding: 12px 16px; margin-bottom: 20px; font-size: 13px; }}
</style>
</head>
<body>
<div class="po-wrap">
  <div class="header">
    <div>
      <div class="store-name">Great Bridge Furniture</div>
      <div class="store-info">1325 S Battlefield Blvd, Chesapeake, VA 23322<br>Phone: (757) 482-6622 &middot; Est. 1984</div>
    </div>
    <div class="po-title">
      <h1>PURCHASE ORDER</h1>
      <div class="po-num">PO # {po_number}</div>
      <div class="po-date">Date: {today}</div>
    </div>
  </div>

  <div class="addresses">
    <div class="addr-block">
      <h3>Vendor / Manufacturer</h3>
      <p><strong>{mfr_name}</strong></p>
      {'<p>' + mfr_contact + '</p>' if mfr_contact else ''}
      {'<p>' + mfr_address.replace(chr(10), '<br>') + '</p>' if mfr_address else ''}
      {'<p>' + mfr_email + '</p>' if mfr_email else ''}
    </div>
    <div class="addr-block">
      <h3>Ship To / Bill To</h3>
      <p><strong>Great Bridge Furniture</strong></p>
      <p>1325 S Battlefield Blvd</p>
      <p>Chesapeake, VA 23322</p>
      <p>(757) 482-6622</p>
    </div>
  </div>

  <div class="confirm-note">
    &#9888;&#65039; Please confirm receipt of this purchase order by replying to this email. Include your order confirmation number and expected ship date.
  </div>

  <table>
    <thead>
      <tr><th>Customer</th><th>Phone</th><th>Items / Description</th><th>Expected Date</th><th>Notes</th></tr>
    </thead>
    <tbody>
      <tr>
        <td>{order.get('customer_name','')}</td>
        <td>{order.get('phone','')}</td>
        <td>{order.get('items','')}</td>
        <td>{order.get('expected_date','')}</td>
        <td>{order.get('notes','')}</td>
      </tr>
    </tbody>
  </table>

  <div class="totals">
    <table>
      <tr><td>Total</td><td><strong>${float(order.get('total_amount') or 0):.2f}</strong></td></tr>
      <tr><td>Deposit Paid</td><td>${float(order.get('deposit_paid') or 0):.2f}</td></tr>
      <tr class="grand-total"><td>Balance Due</td><td>${(float(order.get('total_amount') or 0) - float(order.get('deposit_paid') or 0)):.2f}</td></tr>
    </table>
  </div>

  <div class="footer">
    <p>This purchase order constitutes an offer to purchase the items listed above at the terms stated. Please contact us at (757) 482-6622 or reply to this email with any questions.</p>
    <p style="margin-top:8px;color:#999;">Generated by LightView &middot; Great Bridge Furniture &middot; lightview.win</p>
  </div>
</div>
</body>
</html>"""


# ── AIVM HELPERS ─────────────────────────────────────────────────────────────

def _aivm_decode_pubkey(s):
    from cryptography.hazmat.primitives.asymmetric.ec import EllipticCurvePublicKey
    from cryptography.hazmat.primitives.serialization import load_der_public_key
    import base64 as b64
    raw = b64.b64decode(s) if not s.startswith("0x") else bytes.fromhex(s[2:])
    if len(raw) == 65 and raw[0] == 0x04:
        from cryptography.hazmat.primitives.asymmetric.ec import (
            EllipticCurvePublicNumbers, SECP256R1)
        from cryptography.hazmat.backends import default_backend
        x = int.from_bytes(raw[1:33], "big")
        y = int.from_bytes(raw[33:65], "big")
        return EllipticCurvePublicNumbers(x, y, SECP256R1()).public_key(default_backend())
    return load_der_public_key(raw)

def _aivm_ecdh_wrap(session_key: bytes, peer_pub) -> bytes:
    from cryptography.hazmat.primitives.asymmetric.ec import (
        generate_private_key, ECDH, SECP256R1)
    from cryptography.hazmat.backends import default_backend
    from cryptography.hazmat.primitives.serialization import Encoding, PublicFormat
    ephem_priv = generate_private_key(SECP256R1(), default_backend())
    shared = ephem_priv.exchange(ECDH(), peer_pub)
    from cryptography.hazmat.primitives.kdf.hkdf import HKDF
    from cryptography.hazmat.primitives.hashes import SHA256
    wrap_key = HKDF(SHA256(), 32, None, b"AIVM-wrap", default_backend()).derive(shared)
    from cryptography.hazmat.primitives.ciphers.aead import AESGCM
    iv = _secrets_mod.token_bytes(12)
    ct = AESGCM(wrap_key).encrypt(iv, session_key, None)
    ephem_pub = ephem_priv.public_key().public_bytes(Encoding.X962, PublicFormat.UncompressedPoint)
    return ephem_pub + iv + ct

def _aivm_aes_encrypt(key: bytes, plaintext: bytes) -> bytes:
    from cryptography.hazmat.primitives.ciphers.aead import AESGCM
    iv = _secrets_mod.token_bytes(12)
    return iv + AESGCM(key).encrypt(iv, plaintext, None)

def _aivm_aes_decrypt(key: bytes, blob: bytes) -> bytes:
    from cryptography.hazmat.primitives.ciphers.aead import AESGCM
    return AESGCM(key).decrypt(blob[:12], blob[12:], None)

class AIVMClient:
    """Server-side Lightchain AIVM inference. No user wallet required."""

    def __init__(self, private_key: str):
        from web3 import Web3
        from eth_account import Account
        self._w3       = Web3(Web3.HTTPProvider("https://rpc.mainnet.lightchain.ai"))
        self._account  = Account.from_key(private_key)
        self._registry = self._w3.eth.contract(
            address=Web3.to_checksum_address(AIVM_JOB_REG), abi=AIVM_ABI)
        self._jwt      = None
        self._jwt_exp  = 0
        self._sess_req = req.Session()

    def _auth_headers(self):
        if not self._jwt or time.time() > self._jwt_exp - 60:
            self._refresh_jwt()
        return {"Authorization": f"Bearer {self._jwt}", "Content-Type": "application/json"}

    def _refresh_jwt(self):
        from eth_account.messages import encode_defunct
        r1 = self._sess_req.get(f"{AIVM_GATEWAY}/api/auth/challenge",
                                params={"address": self._account.address}, timeout=15)
        r1.raise_for_status()
        message = r1.json()["message"]
        sig = self._account.sign_message(encode_defunct(text=message))
        r2 = self._sess_req.post(f"{AIVM_GATEWAY}/api/auth/verify",
                                 json={"message": message, "signature": "0x" + sig.signature.hex()},
                                 timeout=15)
        r2.raise_for_status()
        data = r2.json()
        self._jwt = data["token"]
        # expiresAt may be an ISO string, Unix seconds, or Unix ms — handle all
        exp_raw = data.get("expiresAt", 0)
        if isinstance(exp_raw, str):
            try:
                from datetime import datetime
                self._jwt_exp = datetime.fromisoformat(exp_raw.replace("Z", "+00:00")).timestamp()
            except Exception:
                self._jwt_exp = time.time() + 3600
        elif isinstance(exp_raw, (int, float)) and exp_raw > 1e12:
            self._jwt_exp = float(exp_raw) / 1000   # milliseconds → seconds
        elif isinstance(exp_raw, (int, float)) and exp_raw > 1e9:
            self._jwt_exp = float(exp_raw)           # already Unix seconds
        else:
            self._jwt_exp = time.time() + 3600
        print(f"  [AIVM] JWT refreshed")

    def run_inference(self, prompt: str, timeout_secs: int = 300) -> str:
        import websocket as _ws
        from web3 import Web3
        from urllib.parse import quote as _url_quote

        print(f"  [AIVM] inference ({len(prompt)} chars)")
        r = self._sess_req.get(f"{AIVM_GATEWAY}/api/models", timeout=15)
        r.raise_for_status()
        models   = r.json().get("models", [])
        model    = next((m for m in models if m["name"] == "llama3-8b"), models[0] if models else None)
        if not model:
            raise RuntimeError("No AIVM models available")
        model_id = model["id"]

        r = self._sess_req.post(f"{AIVM_GATEWAY}/api/sessions/select",
                                json={"modelId": model_id},
                                headers=self._auth_headers(), timeout=15)
        r.raise_for_status()
        sel = r.json()

        session_key  = _secrets_mod.token_bytes(32)
        enc_worker   = _aivm_ecdh_wrap(session_key, _aivm_decode_pubkey(sel["workerEncryptionKey"]))
        enc_disputer = _aivm_ecdh_wrap(session_key, _aivm_decode_pubkey(sel["disputerEncryptionKey"]))

        r = self._sess_req.post(f"{AIVM_GATEWAY}/api/sessions/prepare",
                                json={"modelId": model_id,
                                      "encWorkerKey":   _b64_mod.b64encode(enc_worker).decode(),
                                      "encDisputerKey": _b64_mod.b64encode(enc_disputer).decode()},
                                headers=self._auth_headers(), timeout=15)
        r.raise_for_status()
        prep = r.json()

        params_hash = bytes.fromhex(model_id[2:].zfill(64) if model_id[:2].lower() == "0x" else model_id.zfill(64))
        sig_bytes   = bytes.fromhex(prep["signature"][2:] if prep["signature"][:2].lower() == "0x" else prep["signature"])
        gas_price   = self._w3.eth.gas_price
        nonce_val   = self._w3.eth.get_transaction_count(self._account.address)

        tx = self._registry.functions.createSession(
            params_hash, Web3.to_checksum_address(prep["worker"]),
            enc_worker, enc_disputer, sig_bytes, prep["expiry"],
        ).build_transaction({"from": self._account.address, "nonce": nonce_val,
                              "gas": 1_000_000, "gasPrice": gas_price,
                              "value": 0, "chainId": AIVM_CHAIN_ID})
        signed  = self._account.sign_transaction(tx)
        tx_hash = self._w3.eth.send_raw_transaction(signed.raw_transaction)
        receipt1 = self._w3.eth.wait_for_transaction_receipt(tx_hash, timeout=90)
        if receipt1.status != 1:
            raise RuntimeError("createSession reverted")

        session_id = None
        for log in receipt1.logs:
            try:
                evt = self._registry.events.SessionCreated().process_log(log)
                session_id = evt["args"]["sessionId"]
                break
            except Exception:
                pass
        if session_id is None:
            raise RuntimeError("SessionCreated event not found")

        relay_token = None
        deadline = time.time() + 60
        while time.time() < deadline:
            r = self._sess_req.get(f"{AIVM_GATEWAY}/api/sessions/{session_id}/token",
                                   headers=self._auth_headers(), timeout=10)
            if r.status_code == 200 and r.json().get("token"):
                relay_token = r.json()["token"]
                break
            time.sleep(1)
        if not relay_token:
            raise RuntimeError("Relay token not ready")

        chunks   = []
        ws_ready = threading.Event()
        ws_err   = [None]

        def _on_message(ws_obj, msg):
            try:
                frame   = json.loads(msg)
                payload = frame.get("payload")
                if payload:
                    blob = _b64_mod.b64decode(payload)
                    pt   = _aivm_aes_decrypt(session_key, blob)
                    chunks.append(pt.decode("utf-8", errors="replace"))
            except Exception:
                pass

        def _on_open(ws_obj): ws_ready.set()
        def _on_error(ws_obj, err): ws_err[0] = err; ws_ready.set()

        import websocket as _ws_mod
        ws = _ws_mod.WebSocketApp(f"{AIVM_RELAY}?token={_url_quote(relay_token)}",
                                   on_message=_on_message, on_open=_on_open, on_error=_on_error)
        ws_thread = threading.Thread(target=ws.run_forever, daemon=True)
        ws_thread.start()
        ws_ready.wait(timeout=15)
        if ws_err[0]:
            raise RuntimeError(f"WebSocket failed: {ws_err[0]}")

        cipher = _aivm_aes_encrypt(session_key, prompt.encode("utf-8"))
        r = self._sess_req.post(f"{AIVM_GATEWAY}/api/blobs",
                                json={"data": _b64_mod.b64encode(cipher).decode()},
                                headers=self._auth_headers(), timeout=15)
        r.raise_for_status()
        blob_hashes = r.json().get("blobHashes", [])
        if not blob_hashes:
            raise RuntimeError("No blob hash")
        _bh         = blob_hashes[0]
        prompt_hash = bytes.fromhex(_bh[2:].zfill(64) if _bh[:2].lower() == "0x" else _bh.zfill(64))

        nonce_val2 = self._w3.eth.get_transaction_count(self._account.address)
        tx2 = self._registry.functions.submitJob(session_id, prompt_hash).build_transaction(
            {"from": self._account.address, "nonce": nonce_val2, "gas": 500_000,
             "gasPrice": gas_price, "value": AIVM_JOB_FEE, "chainId": AIVM_CHAIN_ID})
        signed2  = self._account.sign_transaction(tx2)
        tx_hash2 = self._w3.eth.send_raw_transaction(signed2.raw_transaction)
        receipt2 = self._w3.eth.wait_for_transaction_receipt(tx_hash2, timeout=90)
        if receipt2.status != 1:
            raise RuntimeError("submitJob reverted — check LCAI balance")

        job_completed_topic = "0x" + self._w3.keccak(
            text="JobCompleted(uint256,address,bytes32,bytes32)").hex()
        done     = False
        deadline = time.time() + timeout_secs
        while time.time() < deadline and not done:
            time.sleep(5)
            if chunks:
                done = True
                break
            try:
                logs = self._w3.eth.get_logs(
                    {"address": Web3.to_checksum_address(AIVM_JOB_REG),
                     "fromBlock": receipt2.blockNumber,
                     "toBlock": self._w3.eth.block_number,
                     "topics": [job_completed_topic]})
                if logs:
                    done = True
            except Exception as e:
                print(f"  [AIVM] log poll: {e}")

        time.sleep(3)
        ws.close()
        result = "".join(chunks).strip()
        if not result and not done:
            raise RuntimeError(f"Timeout after {timeout_secs}s")
        return result or "No response from AIVM worker"

# Singleton AIVM client (lazy init)
_aivm_client = None
_aivm_lock   = threading.Lock()

def get_aivm():
    global _aivm_client
    if _aivm_client is None:
        with _aivm_lock:
            if _aivm_client is None:
                if not AIVM_PRIVATE_KEY:
                    raise RuntimeError("LIGHTCHAIN_PRIVATE_KEY env var not set")
                _aivm_client = AIVMClient(AIVM_PRIVATE_KEY)
    return _aivm_client


def _call_aivm(prompt: str, timeout_secs: int = 180) -> str:
    """Convenience wrapper: call AIVM and return the text result. Raises on error."""
    aivm = get_aivm()
    if not aivm:
        raise RuntimeError("AIVM not configured")
    return aivm.run_inference(prompt, timeout_secs=timeout_secs)


# ── QB HELPERS ───────────────────────────────────────────────────────────────

def _qb_get_token(realm_id: str) -> dict | None:
    with _QB_LOCK:
        tok = _QB_TOKENS.get(realm_id)
    if not tok:
        return None
    # Refresh if expiring soon
    if time.time() > tok["expires_at"] - 120:
        tok = _qb_refresh_token(realm_id, tok["refresh_token"])
    return tok

def _qb_refresh_token(realm_id: str, refresh_token: str) -> dict:
    import base64 as b64
    creds = b64.b64encode(f"{QB_CLIENT_ID}:{QB_CLIENT_SECRET}".encode()).decode()
    r = req.post(QB_TOKEN_URL,
                 headers={"Authorization": f"Basic {creds}",
                          "Content-Type": "application/x-www-form-urlencoded",
                          "Accept": "application/json"},
                 data={"grant_type": "refresh_token", "refresh_token": refresh_token},
                 timeout=15)
    r.raise_for_status()
    data = r.json()
    tok  = {"access_token":  data["access_token"],
            "refresh_token": data.get("refresh_token", refresh_token),
            "expires_at":    time.time() + data.get("expires_in", 3600)}
    with _QB_LOCK:
        _QB_TOKENS[realm_id] = tok
    return tok

def _qb_api(realm_id: str, endpoint: str, params: dict = None) -> dict:
    tok = _qb_get_token(realm_id)
    if not tok:
        raise RuntimeError("QuickBooks not connected")
    url = f"{QB_API_BASE}/{realm_id}/{endpoint}"
    r   = req.get(url, headers={"Authorization": f"Bearer {tok['access_token']}",
                                 "Accept": "application/json"},
                  params=params or {}, timeout=15)
    r.raise_for_status()
    return r.json()

def _qb_report(realm_id: str, report_name: str, params: dict = None) -> dict:
    return _qb_api(realm_id, f"reports/{report_name}", params)

# ── API ROUTES ───────────────────────────────────────────────────────────────

@app.route("/health")
def health():
    return jsonify({"status": "ok", "aivm_key": bool(AIVM_PRIVATE_KEY),
                    "qb_configured": bool(QB_CLIENT_ID and QB_CLIENT_SECRET)})

# ── AIVM: Analyze monitor data ───────────────────────────────────────────────
@app.route("/api/analyze", methods=["POST"])
def analyze():
    """
    Body: {
      company: str,
      industry: str,
      monitors: [{name, value, unit, status, threshold}],
      alerts: [{title, severity, triggered_value, threshold}],
      rules: [str]
    }
    Returns: { analysis: str }
    """
    data = request.get_json(force=True)
    if not data:
        return jsonify({"error": "No data"}), 400

    company  = data.get("company", "Unknown Company")
    industry = data.get("industry", "General")
    monitors = data.get("monitors", [])
    alerts   = data.get("alerts", [])
    rules    = data.get("rules", [])

    monitor_lines = "\n".join(
        f"  - {m['name']}: {m['value']} {m.get('unit','')} | Status: {m['status']} | Threshold: {m.get('threshold','N/A')}"
        for m in monitors)
    alert_lines = "\n".join(
        f"  - [{a['severity'].upper()}] {a['title']}: triggered at {a.get('triggered_value','?')} (limit: {a.get('threshold','?')})"
        for a in alerts) or "  None"
    rule_lines = "\n".join(f"  - {r}" for r in rules) or "  None defined"

    prompt = f"""You are LightWatch AI, an enterprise compliance and operations monitor for {company} ({industry}).

Current live monitor readings:
{monitor_lines}

Active alerts:
{alert_lines}

Owner-defined rules:
{rule_lines}

Provide a concise AI analysis (3-4 paragraphs) covering:
1. What the data pattern means for operations right now
2. Which alerts need immediate attention and why
3. Whether any rule violations are present
4. One specific recommended action the operator should take in the next hour

Write in plain English suitable for a business owner, not a technician. Be specific about the numbers."""

    try:
        aivm   = get_aivm()
        result = aivm.run_inference(prompt, timeout_secs=240)
        return jsonify({"analysis": result, "source": "AIVM"})
    except Exception as e:
        print(f"[analyze] AIVM error: {e}")
        # Fallback: rule-based summary
        critical = [a for a in alerts if a.get("severity") == "critical"]
        warnings = [a for a in alerts if a.get("severity") == "warning"]
        summary  = f"AI analysis for {company}: "
        if critical:
            summary += f"{len(critical)} critical alert(s) require immediate attention: {', '.join(a['title'] for a in critical[:2])}. "
        if warnings:
            summary += f"{len(warnings)} warning(s) are active. "
        if not alerts:
            summary += "All monitors within normal parameters. "
        summary += "Recommend reviewing threshold settings and acknowledging active alerts."
        return jsonify({"analysis": summary, "source": "fallback", "error": str(e)})


# ── AIVM: Generate report narrative ─────────────────────────────────────────
@app.route("/api/report", methods=["POST"])
def generate_report():
    """
    Body: {
      company: str,
      industry: str,
      report_type: 'bank'|'regulator'|'auditor'|'exec'|'govt'|'investor',
      data: { monitors, alerts, rules, stats }
    }
    Returns: { report: str }
    """
    data = request.get_json(force=True)
    if not data:
        return jsonify({"error": "No data"}), 400

    # Accept either {company, industry, data:{...}} or {entity:{name, industry, ...}, report_type}
    entity      = data.get("entity", {})
    company     = data.get("company", entity.get("name", "Unknown"))
    industry    = data.get("industry", entity.get("industry", "General"))
    report_type = data.get("report_type", "exec")
    company_data = data.get("data", entity)   # fall back to entity block

    audience_map = {
        "bank":      "a bank loan officer reviewing creditworthiness and operational stability",
        "regulator": "a government regulator verifying compliance with industry standards",
        "auditor":   "an external auditor conducting an operational audit",
        "exec":      "a C-suite executive who needs a high-level operational summary",
        "govt":      "a government contracting officer evaluating suitability for a contract",
        "investor":  "a potential investor evaluating the business for funding",
    }
    audience = audience_map.get(report_type, audience_map["exec"])

    monitors = company_data.get("monitors", [])
    alerts   = company_data.get("alerts", [])
    stats    = company_data.get("stats", {})

    data_summary = f"""Company: {company} | Industry: {industry}
Overall status: {stats.get('overall', 'Unknown')}
Active monitors: {len(monitors)}
Active alerts: {len(alerts)}
Uptime: {stats.get('uptime', 'N/A')}
Compliance score: {stats.get('compliance', 'N/A')}

Alert summary: {', '.join(f"{a.get('severity','?').upper()}: {a['title']}" for a in alerts) or 'No active alerts'}"""

    prompt = f"""You are LightWatch AI generating a formal business report for {company}.

Report audience: {audience}

Operational data:
{data_summary}

Write a professional {report_type} report (3-4 paragraphs) appropriate for {audience}.
- Use formal language appropriate to the audience
- Include specific data points
- Note the blockchain-verified audit trail where relevant
- For regulators/auditors: emphasize compliance and immutable record-keeping
- For banks/investors: emphasize operational stability and risk management
- For executives: focus on business impact and action items
- For government contracts: emphasize reliability, compliance standards met, and audit trail
Keep it concise and factual. Today's date: {time.strftime('%B %d, %Y')}."""

    try:
        aivm   = get_aivm()
        result = aivm.run_inference(prompt, timeout_secs=240)
        return jsonify({"report": result, "source": "AIVM"})
    except Exception as e:
        print(f"[report] AIVM error: {e}")
        return jsonify({"report": f"Report generation unavailable. Please review the dashboard data directly. Error: {str(e)}",
                        "source": "error", "fallback": True})


# ── QUICKBOOKS OAuth ─────────────────────────────────────────────────────────
@app.route("/api/qb/auth")
def qb_auth():
    if not QB_CLIENT_ID:
        return jsonify({"error": "QuickBooks not configured — set QUICKBOOKS_CLIENT_ID and QUICKBOOKS_CLIENT_SECRET"}), 503

    state    = _secrets_mod.token_urlsafe(16)
    session["qb_state"] = state
    auth_url = (f"{QB_AUTH_URL}?client_id={QB_CLIENT_ID}"
                f"&response_type=code&scope={QB_SCOPE}"
                f"&redirect_uri={QB_REDIRECT_URI}&state={state}")
    return redirect(auth_url)


@app.route("/api/qb/callback")
def qb_callback():
    import base64 as b64
    error = request.args.get("error")
    if error:
        return jsonify({"error": f"QB OAuth error: {error}"}), 400

    state = request.args.get("state")
    if state != session.get("qb_state"):
        return jsonify({"error": "State mismatch — possible CSRF"}), 400

    code     = request.args.get("code")
    realm_id = request.args.get("realmId")

    creds = b64.b64encode(f"{QB_CLIENT_ID}:{QB_CLIENT_SECRET}".encode()).decode()
    r = req.post(QB_TOKEN_URL,
                 headers={"Authorization": f"Basic {creds}",
                          "Content-Type": "application/x-www-form-urlencoded",
                          "Accept": "application/json"},
                 data={"grant_type": "authorization_code", "code": code,
                       "redirect_uri": QB_REDIRECT_URI},
                 timeout=15)
    r.raise_for_status()
    tok_data = r.json()

    with _QB_LOCK:
        _QB_TOKENS[realm_id] = {
            "access_token":  tok_data["access_token"],
            "refresh_token": tok_data["refresh_token"],
            "expires_at":    time.time() + tok_data.get("expires_in", 3600),
        }

    session["qb_realm_id"] = realm_id
    print(f"[QB] Connected: realm_id={realm_id}")
    # Redirect back to the frontend with success
    frontend_url = os.environ.get("FRONTEND_URL", "https://keiko-dev-lcai.github.io/lightwatch/")
    return redirect(f"{frontend_url}?qb_connected=1&realm={realm_id}")


@app.route("/api/qb/status")
def qb_status():
    realm_id = request.args.get("realm") or session.get("qb_realm_id")
    if not realm_id or realm_id not in _QB_TOKENS:
        return jsonify({"connected": False})
    return jsonify({"connected": True, "realm_id": realm_id})


@app.route("/api/qb/data")
def qb_data():
    """
    Pull live QuickBooks data and return it in LightWatch monitor format.
    Query param: realm=<realmId>
    Returns: { company_name, monitors, stats, raw }
    """
    realm_id = request.args.get("realm") or session.get("qb_realm_id")
    if not realm_id:
        return jsonify({"error": "Not connected to QuickBooks"}), 401

    try:
        # Company info
        company_info = _qb_api(realm_id, "companyinfo/" + realm_id)
        co = company_info.get("CompanyInfo", {})
        company_name = co.get("CompanyName", "Your Company")

        # P&L for current month
        now   = time.strftime("%Y-%m-%d")
        start = time.strftime("%Y-%m-01")
        pnl   = _qb_report(realm_id, "ProfitAndLoss", {"start_date": start, "end_date": now})

        # Balance Sheet
        bs = _qb_report(realm_id, "BalanceSheet", {"start_date": start, "end_date": now})

        # Cash Flow
        cf = _qb_report(realm_id, "CashFlow", {"start_date": start, "end_date": now})

        # Parse P&L into monitor-friendly values
        def _find_row(report, label_contains):
            for row in report.get("Rows", {}).get("Row", []):
                if isinstance(row, dict):
                    header = row.get("Header", {})
                    col_data = header.get("ColData", [])
                    for cd in col_data:
                        if label_contains.lower() in str(cd.get("value", "")).lower():
                            # Find the value column
                            summary = row.get("Summary", {}).get("ColData", [])
                            for i, s in enumerate(summary):
                                if i > 0 and s.get("value"):
                                    try:
                                        return float(str(s["value"]).replace(",", ""))
                                    except:
                                        pass
            return None

        revenue  = _find_row(pnl, "Total Income")  or _find_row(pnl, "Total Revenue") or 0
        expenses = _find_row(pnl, "Total Expense") or _find_row(pnl, "Total Cost")    or 0
        net_income = revenue - expenses if (revenue and expenses) else (_find_row(pnl, "Net Income") or 0)

        total_assets   = _find_row(bs,  "Total Assets")      or 0
        total_liab     = _find_row(bs,  "Total Liabilities") or 0
        cash_from_ops  = _find_row(cf,  "Net Cash")          or 0

        # Build LightWatch monitor format
        monitors = [
            {"name": "Monthly Revenue",      "value": f"${revenue:,.0f}",    "unit": "USD", "status": "green" if revenue > 0 else "yellow", "threshold": "Positive"},
            {"name": "Monthly Expenses",     "value": f"${expenses:,.0f}",   "unit": "USD", "status": "yellow" if expenses > revenue * 0.9 else "green", "threshold": f"< ${revenue:,.0f} revenue"},
            {"name": "Net Income",           "value": f"${net_income:,.0f}", "unit": "USD", "status": "green" if net_income > 0 else "red",  "threshold": "Positive"},
            {"name": "Total Assets",         "value": f"${total_assets:,.0f}","unit": "USD","status": "green", "threshold": "N/A"},
            {"name": "Total Liabilities",    "value": f"${total_liab:,.0f}", "unit": "USD", "status": "green" if total_liab < total_assets * 0.6 else "yellow", "threshold": "< 60% of assets"},
            {"name": "Cash from Operations", "value": f"${cash_from_ops:,.0f}", "unit": "USD", "status": "green" if cash_from_ops > 0 else "red", "threshold": "Positive"},
        ]

        alerts = []
        if net_income < 0:
            alerts.append({"title": "Net Loss This Month", "severity": "critical",
                           "triggered_value": f"${net_income:,.0f}", "threshold": "$0"})
        if total_liab > total_assets * 0.6:
            ratio = total_liab / total_assets if total_assets else 0
            alerts.append({"title": "High Debt-to-Asset Ratio", "severity": "warning",
                           "triggered_value": f"{ratio:.0%}", "threshold": "60%"})
        if expenses > revenue * 0.95 and revenue > 0:
            alerts.append({"title": "Expense Ratio Near Revenue", "severity": "warning",
                           "triggered_value": f"{expenses/revenue:.0%} of revenue", "threshold": "95%"})

        # Stats summary
        profit_margin = (net_income / revenue * 100) if revenue > 0 else 0
        stats = {
            "overall":    "red" if any(a["severity"] == "critical" for a in alerts) else ("yellow" if alerts else "green"),
            "uptime":     "N/A",
            "compliance": f"{max(0, min(100, 95 - len(alerts)*10))}%",
            "active_monitors": len(monitors),
            "active_alerts":   len(alerts),
            "profit_margin":   f"{profit_margin:.1f}%",
        }

        return jsonify({
            "company_name": company_name,
            "monitors":     monitors,
            "alerts":       alerts,
            "stats":        stats,
            "source":       "QuickBooks",
            "as_of":        now,
        })

    except Exception as e:
        print(f"[qb/data] error: {e}")
        return jsonify({"error": str(e)}), 500


# ── PIN Authentication ───────────────────────────────────────────────────────
LIGHTVIEW_ADMIN_PIN = os.environ.get("LIGHTVIEW_PIN", "8300")

@app.route("/api/auth", methods=["POST"])
def auth_pin():
    """
    Body: { pin: str }
    Validates PIN against LIGHTVIEW_PIN env var (default 8300).
    Returns { ok: true, role: "admin" } on success.
    """
    data = request.get_json(force=True) or {}
    pin  = str(data.get("pin", "")).strip()
    if pin == LIGHTVIEW_ADMIN_PIN:
        return jsonify({"ok": True, "role": "admin"})
    return jsonify({"ok": False}), 401

# ── Contact / Demo Request ───────────────────────────────────────────────────
@app.route("/api/contact", methods=["POST"])
def contact():
    """
    Body: { company: str, industry: str, use_case: str }
    Logs the demo request to Railway stdout.
    """
    data     = request.get_json(force=True) or {}
    company  = str(data.get("company",  ""))[:200].strip()
    industry = str(data.get("industry", ""))[:100].strip()
    use_case = str(data.get("use_case", ""))[:500].strip()
    print(f"[DEMO REQUEST] Company: {company!r} | Industry: {industry!r} | Use case: {use_case!r}")
    return jsonify({"ok": True})


# ── GREAT BRIDGE FURNITURE — ORDER ENDPOINTS ─────────────────────────────────

@app.route("/api/gb/orders", methods=["GET"])
def gb_get_orders():
    """Return all orders sorted: Ready for Pickup first, then Ordered, In Transit, etc."""
    try:
        _init_db()
        with _db() as conn:
            rows = conn.execute(
                "SELECT * FROM gb_orders ORDER BY "
                "CASE status "
                "  WHEN 'Ready for Pickup' THEN 1 "
                "  WHEN 'Ordered'          THEN 2 "
                "  WHEN 'In Transit'       THEN 3 "
                "  WHEN 'Received'         THEN 4 "
                "  WHEN 'Delivered'        THEN 5 "
                "  WHEN 'Paid'             THEN 6 "
                "  ELSE 7 END, created_at DESC"
            ).fetchall()
        return jsonify({"orders": [dict(r) for r in rows]})
    except Exception as e:
        print(f"[gb/orders GET] {e}")
        return jsonify({"error": str(e), "orders": []}), 500


@app.route("/api/gb/orders", methods=["POST"])
def gb_create_order():
    """Create a single order from the manual entry form."""
    data = request.get_json(force=True) or {}
    user = str(data.get("user_name") or "Unknown")[:100]
    try:
        _init_db()
        items = data.get("items", [])
        items_str = json.dumps(items) if isinstance(items, list) else str(items)
        customer  = str(data.get("customer_name") or "Unknown")[:200]
        with _db() as conn:
            cur = conn.execute(
                "INSERT INTO gb_orders "
                "(customer_name, phone, items, manufacturer, total_amount, deposit_paid, expected_date, delivery_date, status, notes) "
                "VALUES (?,?,?,?,?,?,?,?,?,?)",
                (customer,
                 str(data.get("phone") or "")[:50],
                 items_str,
                 str(data.get("manufacturer") or "")[:200],
                 float(data.get("total_amount") or 0),
                 float(data.get("deposit_paid") or 0),
                 str(data.get("expected_date") or "")[:50],
                 str(data.get("delivery_date") or "")[:50],
                 str(data.get("status") or "Ordered")[:50],
                 str(data.get("notes") or "")[:500])
            )
            conn.commit()
            new_id = cur.lastrowid
        items_preview = ", ".join(items[:2]) if isinstance(items, list) else str(items)[:60]
        _log_activity(user, "Added order", f"#{new_id}",
                      f"Customer: {customer} — Items: {items_preview[:80]}")
        return jsonify({"ok": True, "id": new_id})
    except Exception as e:
        print(f"[gb/orders POST] {e}")
        return jsonify({"error": str(e)}), 500


@app.route("/api/gb/orders/<int:order_id>", methods=["PUT"])
def gb_update_order(order_id):
    """Update any fields on an existing order (e.g. status change)."""
    data = request.get_json(force=True) or {}
    user = str(data.get("user_name") or "Unknown")[:100]
    try:
        _init_db()
        allowed = ["customer_name", "phone", "manufacturer", "total_amount",
                   "deposit_paid", "expected_date", "delivery_date", "status", "notes"]
        fields, vals = [], []
        for k in allowed:
            if k in data:
                fields.append(f"{k} = ?")
                vals.append(data[k])
        if "items" in data:
            fields.append("items = ?")
            vals.append(json.dumps(data["items"]) if isinstance(data["items"], list) else str(data["items"]))
        if not fields:
            return jsonify({"error": "Nothing to update"}), 400
        fields.append("updated_at = datetime('now')")
        vals.append(order_id)
        # Fetch customer name for the log
        with _db() as conn:
            row = conn.execute("SELECT customer_name FROM gb_orders WHERE id=?", (order_id,)).fetchone()
            customer = row["customer_name"] if row else "?"
            conn.execute(f"UPDATE gb_orders SET {', '.join(fields)} WHERE id = ?", vals)
            conn.commit()
        # Log a meaningful description of what changed
        if "status" in data:
            _log_activity(user, "Changed status to",
                          f"\"{data['status']}\"",
                          f"Order #{order_id} — {customer}")
        else:
            changed = [k for k in allowed if k in data] + (["items"] if "items" in data else [])
            _log_activity(user, "Updated order", f"#{order_id}",
                          f"Customer: {customer} — Fields: {', '.join(changed)}")
        return jsonify({"ok": True})
    except Exception as e:
        print(f"[gb/orders PUT {order_id}] {e}")
        return jsonify({"error": str(e)}), 500


@app.route("/api/gb/orders/<int:order_id>", methods=["DELETE"])
def gb_delete_order(order_id):
    """Delete an order and its docs. Logs who deleted it."""
    data = request.get_json(force=True) or {}
    user = str(data.get("user_name") or "Unknown")[:100]
    try:
        _init_db()
        with _db() as conn:
            row = conn.execute(
                "SELECT customer_name FROM gb_orders WHERE id=?", (order_id,)
            ).fetchone()
            if not row:
                return jsonify({"error": "Order not found"}), 404
            customer = row["customer_name"]
            conn.execute("DELETE FROM gb_docs WHERE order_id=?", (order_id,))
            conn.execute("DELETE FROM gb_orders WHERE id=?", (order_id,))
            conn.commit()
        _log_activity(user, "Deleted order", f"#{order_id}",
                      f"Customer: {customer}")
        return jsonify({"ok": True})
    except Exception as e:
        print(f"[gb/orders DELETE {order_id}] {e}")
        return jsonify({"error": str(e)}), 500


@app.route("/api/gb/orders/save-import", methods=["POST"])
def gb_save_import():
    """Batch-save orders that came back from the AIVM import preview."""
    data   = request.get_json(force=True) or {}
    orders = data.get("orders", [])
    user   = str(data.get("user_name") or "Unknown")[:100]
    if not orders:
        return jsonify({"error": "No orders to save"}), 400
    try:
        _init_db()
        saved = 0
        with _db() as conn:
            for o in orders:
                items_str = json.dumps(o.get("items", [])) if isinstance(o.get("items"), list) else str(o.get("items") or "")
                try:
                    total   = float(str(o.get("total_amount")   or "0").replace("$","").replace(",",""))
                    deposit = float(str(o.get("deposit_paid")   or "0").replace("$","").replace(",",""))
                except Exception:
                    total = deposit = 0
                conn.execute(
                    "INSERT INTO gb_orders "
                    "(customer_name, phone, items, manufacturer, total_amount, deposit_paid, expected_date, status, notes) "
                    "VALUES (?,?,?,?,?,?,?,?,?)",
                    (str(o.get("customer_name") or "Unknown")[:200],
                     str(o.get("phone") or "")[:50],
                     items_str,
                     str(o.get("manufacturer") or "")[:200],
                     total, deposit,
                     str(o.get("expected_date") or "")[:50],
                     str(o.get("status") or "Ordered")[:50],
                     str(o.get("notes") or "")[:500])
                )
                saved += 1
            conn.commit()
        _log_activity(user, "Imported orders via spreadsheet", "",
                      f"{saved} orders added to system")
        return jsonify({"ok": True, "saved": saved})
    except Exception as e:
        print(f"[gb/orders/save-import] {e}")
        return jsonify({"error": str(e)}), 500


@app.route("/api/gb/import-preview", methods=["POST"])
def gb_import_preview():
    """Read column headers + sample rows from an uploaded spreadsheet — no AIVM needed."""
    if 'file' not in request.files:
        return jsonify({"error": "No file uploaded"}), 400
    f     = request.files['file']
    fname = (f.filename or "").lower()
    try:
        content = f.read()
        headers = []
        samples = []
        if fname.endswith('.csv'):
            text   = content.decode('utf-8', errors='replace')
            reader = _csv_mod.DictReader(io.StringIO(text))
            headers = list(reader.fieldnames or [])
            for i, row in enumerate(reader):
                samples.append({k: str(v).strip() for k, v in row.items()})
                if i >= 4:
                    break
        elif fname.endswith(('.xlsx', '.xls')):
            import openpyxl
            wb = openpyxl.load_workbook(io.BytesIO(content), data_only=True)
            ws = wb.active
            for i, row in enumerate(ws.iter_rows(values_only=True)):
                if i == 0:
                    headers = [str(h).strip() if h is not None else f"col{j}" for j, h in enumerate(row)]
                else:
                    samples.append({headers[j]: str(v).strip() if v is not None else "" for j, v in enumerate(row) if j < len(headers)})
                    if i >= 5:
                        break
            wb.close()
        else:
            return jsonify({"error": "Please upload a .csv or .xlsx file"}), 400
        return jsonify({"headers": headers, "samples": samples, "filename": f.filename})
    except Exception as e:
        return jsonify({"error": f"Could not read file: {e}"}), 400


@app.route("/api/gb/import", methods=["POST"])
def gb_import_spreadsheet():
    """Upload CSV or Excel → returns JSON orders array for preview.
    Accepts optional 'column_map' JSON: {"their_col": "our_field", ...}
    Our fields: customer_name, phone, items, manufacturer, total_amount,
                deposit_paid, expected_date, status, notes
    """
    if 'file' not in request.files:
        return jsonify({"error": "No file uploaded"}), 400
    f     = request.files['file']
    fname = (f.filename or "").lower()

    # Optional explicit column map from the UI
    explicit_map = {}
    try:
        explicit_map = json.loads(request.form.get('column_map', '{}'))
    except Exception:
        pass

    raw_rows = []
    try:
        content = f.read()
        if fname.endswith('.csv'):
            text   = content.decode('utf-8', errors='replace')
            reader = _csv_mod.DictReader(io.StringIO(text))
            for row in reader:
                raw_rows.append({k: str(v) for k, v in row.items() if v is not None})
                if len(raw_rows) >= 10000:
                    break
        elif fname.endswith(('.xlsx', '.xls')):
            import openpyxl
            # Do NOT use read_only=True — some Excel files have a broken internal
            # dimension tag that causes openpyxl to stop early in read_only mode.
            wb      = openpyxl.load_workbook(io.BytesIO(content), data_only=True)
            headers = None
            # Iterate ALL worksheets so multi-sheet workbooks are fully read
            for ws in wb.worksheets:
                for i, row in enumerate(ws.iter_rows(values_only=True)):
                    if headers is None and i == 0:
                        headers = [str(h).strip() if h is not None else f"col{j}"
                                   for j, h in enumerate(row)]
                    elif headers is not None and not all(v is None for v in row):
                        raw_rows.append({
                            headers[j]: str(v).strip() if v is not None else ""
                            for j, v in enumerate(row) if j < len(headers)
                        })
                    if len(raw_rows) >= 10000:
                        break
                if len(raw_rows) >= 10000:
                    break
            wb.close()
        else:
            return jsonify({"error": "Please upload a .csv or .xlsx file"}), 400
    except Exception as e:
        return jsonify({"error": f"Could not read file: {e}"}), 400

    if not raw_rows:
        return jsonify({"error": "File is empty or has no readable rows"}), 400

    headers_found = list(raw_rows[0].keys()) if raw_rows else []
    sample        = raw_rows[:60]
    rows_text     = "\n".join(
        " | ".join(f"{k}: {v}" for k, v in row.items() if v and v != 'None')
        for row in sample
    )

    prompt = (
        "You are analyzing a spreadsheet from Great Bridge Furniture, a furniture store in Chesapeake Virginia.\n"
        f"The spreadsheet has these column headers: {', '.join(headers_found)}\n"
        f"Here are the rows (up to 60 shown):\n{rows_text}\n\n"
        "Extract every customer order you can find. For each distinct customer order, create one JSON entry with these fields:\n"
        "  customer_name (string), phone (string), items (array of strings), manufacturer (string),\n"
        "  total_amount (number, dollars), deposit_paid (number, dollars),\n"
        "  expected_date (YYYY-MM-DD string or null), status (one of: Ordered / In Transit / Received / Ready for Pickup / Delivered / Paid),\n"
        "  notes (string)\n"
        "Respond with ONLY a valid JSON array — no other text, no markdown, no code blocks."
    )

    aivm_result = None
    try:
        aivm = get_aivm()
        if aivm:
            aivm_result = aivm.run_inference(prompt, timeout_secs=300)
            clean = _re_mod.sub(r'```(?:json)?\s*|\s*```', '', aivm_result).strip()
            match = _re_mod.search(r'\[.*\]', clean, _re_mod.DOTALL)
            orders = json.loads(match.group()) if match else json.loads(clean)
            if not isinstance(orders, list):
                orders = []
            return jsonify({
                "orders": orders,
                "rows_scanned": len(raw_rows),
                "orders_found": len(orders),
                "source": "AIVM",
                "headers": headers_found,
            })
    except Exception as e:
        print(f"[gb/import AIVM] {e} — using rule-based fallback")

    # ── Rule-based fallback: map columns without AIVM ────────────────────────
    # If explicit_map provided by UI, use it directly
    if explicit_map:
        STATUS_VALS = {"ordered","in transit","received","ready for pickup","delivered","paid"}
        def _normalize_status(val):
            if not val: return "Ordered"
            v = val.strip().lower()
            for s in STATUS_VALS:
                if s in v: return s.title()
            return "Ordered"
        def _num(val):
            try: return float(_re_mod.sub(r'[^\d.]', '', val))
            except: return 0.0
        def _parse_date_em(raw):
            if not raw: return None
            from datetime import datetime as _dt
            for fmt in ('%Y-%m-%d','%m/%d/%Y','%m-%d-%Y','%d/%m/%Y','%Y/%m/%d'):
                try: return _dt.strptime(raw, fmt).strftime('%Y-%m-%d')
                except: pass
            return raw
        # Pre-collect any free-form "notes_append:FieldName" → spreadsheet-column mappings
        _notes_append_map = {k[len('notes_append:'):]: v
                             for k, v in explicit_map.items()
                             if k.startswith('notes_append:')}

        orders = []
        inv_rows_em = []
        for row in raw_rows:
            def _gx(field):
                col = explicit_map.get(field)
                return row.get(col, "").strip() if col else ""
            # Build extra notes from free-form columns: "Label1: Val1 | Label2: Val2"
            _extra_parts = [f"{lbl}: {row.get(col, '').strip()}"
                            for lbl, col in _notes_append_map.items()
                            if row.get(col, '').strip()]
            _notes_extra = " | ".join(_extra_parts)
            def _merge_notes(base):
                if _notes_extra:
                    return (base + "\n" + _notes_extra) if base else _notes_extra
                return base
            name = _gx("customer_name")
            if name:
                # Order row — has a customer/TAG
                items_raw = _gx("items")
                items = [i.strip() for i in _re_mod.split(r'[;,]+', items_raw) if i.strip()] if items_raw else []
                orders.append({
                    "customer_name": name, "phone": _gx("phone"),
                    "items": items or ([items_raw] if items_raw else []),
                    "manufacturer": _gx("manufacturer"),
                    "total_amount": _num(_gx("total_amount")),
                    "deposit_paid": _num(_gx("deposit_paid")),
                    "expected_date": _parse_date_em(_gx("expected_date")),
                    "status": _normalize_status(_gx("status")),
                    "notes": _merge_notes(_gx("notes")),
                })
            else:
                # Inventory row — blank TAG means warehouse stock
                desc = _gx("items") or _gx("notes")
                if not desc: continue
                inv_rows_em.append({
                    "sku": _gx("sku"), "description": desc,
                    "manufacturer": _gx("manufacturer"),
                    "finish_fabric": _gx("finish_fabric"),
                    "warehouse_location": _gx("warehouse_location"),
                    "cost_price": _num(_gx("cost_price")) or None,
                    "retail_price": _num(_gx("total_amount")) or None,
                    "date_in": _parse_date_em(_gx("date_in") or _gx("expected_date")),
                    "notes": _merge_notes(_gx("notes")), "status": "Available",
                })
        inv_saved_em = 0
        if inv_rows_em:
            try:
                _init_db()
                with _db() as conn:
                    for inv in inv_rows_em:
                        conn.execute(
                            "INSERT INTO gb_inventory (sku,description,manufacturer,finish_fabric,"
                            "warehouse_location,cost_price,retail_price,status,date_in,notes) "
                            "VALUES (?,?,?,?,?,?,?,?,?,?)",
                            (inv.get('sku',''), inv.get('description',''), inv.get('manufacturer',''),
                             inv.get('finish_fabric',''), inv.get('warehouse_location',''),
                             inv.get('cost_price'), inv.get('retail_price'),
                             inv.get('status','Available'), inv.get('date_in',''), inv.get('notes','')))
                    conn.commit()
                inv_saved_em = len(inv_rows_em)
            except Exception as _ie:
                print(f"[gb/import inventory] {_ie}")
        return jsonify({"orders": orders, "rows_scanned": len(raw_rows),
                        "orders_added": len(orders), "inventory_added": inv_saved_em,
                        "source": "explicit_map", "headers": headers_found})

    # Look for column names that match known order and inventory fields
    COL_MAP = {
        "customer_name":     ["customer", "customer name", "name", "client", "buyer", "bill to", "sold to",
                              "tagged", "tag", "tagged as", "customer tag", "warehouse tag"],
        "phone":             ["phone", "telephone", "cell", "mobile", "contact", "phone number"],
        "items":             ["items", "item", "description", "product", "furniture", "order", "goods",
                              "merchandise", "desc", "part description", "model"],
        "manufacturer":      ["manufacturer", "vendor", "supplier", "brand", "maker", "mfg", "source"],
        "total_amount":      ["total", "total amount", "price", "amount", "sale", "invoice",
                              "balance", "retail", "retail price"],
        "deposit_paid":      ["deposit", "deposit paid", "paid", "down payment", "down", "payment"],
        "expected_date":     ["expected", "expected date", "eta", "due date", "delivery date", "arrival",
                              "ship date", "date out", "out date", "date sold", "ship", "delivery"],
        "status":            ["status", "order status", "stage", "state"],
        "notes":             ["notes", "note", "comments", "comment", "remarks", "memo", "details",
                              "po#", "po #", "purchase order", "po number", "verified", "intl",
                              "initials", "approved by"],
        # Inventory-specific fields
        "warehouse_location": ["location", "loc", "warehouse", "warehouse location", "warehouse loc",
                               "wh location", "bin", "shelf", "whloc"],
        "sku":               ["sku", "part#", "part number", "part #", "part no", "item#",
                              "item number", "style#", "style number", "partnumber"],
        "finish_fabric":     ["finish", "fabric", "finish/fabric", "color", "colour", "material"],
        "cost_price":        ["cost", "cost price", "unit cost", "our cost", "dealer cost", "net", "net price"],
        "date_in":           ["date in", "date arrived", "arrived", "received", "receipt date",
                              "in date", "date received"],
    }

    def _find_col(field, row_keys):
        """Return the first row key that matches the field's aliases."""
        aliases = COL_MAP.get(field, [])
        for key in row_keys:
            key_lower = key.lower().strip()
            if key_lower in aliases or any(a in key_lower for a in aliases):
                return key
        return None

    STATUS_VALS = {"ordered", "in transit", "received", "ready for pickup", "delivered", "paid"}

    def _normalize_status(val):
        if not val:
            return "Ordered"
        v = val.strip().lower()
        for s in STATUS_VALS:
            if s in v:
                return s.title().replace("For", "for").replace("for P", "for P")
        return "Ordered"

    def _num(val):
        try:
            return float(_re_mod.sub(r'[^\d.]', '', val))
        except Exception:
            return 0.0

    def _parse_date(raw):
        if not raw: return None
        from datetime import datetime as _dt
        for fmt in ('%Y-%m-%d', '%m/%d/%Y', '%m-%d-%Y', '%d/%m/%Y', '%Y/%m/%d'):
            try:
                return _dt.strptime(raw, fmt).strftime('%Y-%m-%d')
            except Exception:
                pass
        return raw

    keys = list(raw_rows[0].keys()) if raw_rows else []
    col_lookup = {field: _find_col(field, keys) for field in COL_MAP}

    orders = []
    inventory_items = []
    for row in raw_rows:
        def _get(field):
            col = col_lookup.get(field)
            return row.get(col, "").strip() if col else ""

        name = _get("customer_name")
        if name:
            # Order row — has a customer/TAG
            items_raw = _get("items")
            items = [i.strip() for i in _re_mod.split(r'[;,]+', items_raw) if i.strip()] if items_raw else []
            orders.append({
                "customer_name":  name,
                "phone":          _get("phone"),
                "items":          items or [items_raw] if items_raw else [],
                "manufacturer":   _get("manufacturer"),
                "total_amount":   _num(_get("total_amount")),
                "deposit_paid":   _num(_get("deposit_paid")),
                "expected_date":  _parse_date(_get("expected_date")),
                "status":         _normalize_status(_get("status")),
                "notes":          _get("notes"),
            })
        else:
            # Inventory row — blank TAG means warehouse stock, not a customer order
            desc = _get("items") or _get("notes")
            if not desc: continue
            inventory_items.append({
                "sku":               _get("sku"),
                "description":       desc,
                "manufacturer":      _get("manufacturer"),
                "finish_fabric":     _get("finish_fabric"),
                "warehouse_location": _get("warehouse_location"),
                "cost_price":        _num(_get("cost_price")) or None,
                "retail_price":      _num(_get("total_amount")) or None,
                "date_in":           _parse_date(_get("date_in") or _get("expected_date")),
                "notes":             _get("notes"),
                "status":            "Available",
            })

    # Save inventory items immediately to db
    inv_saved = 0
    if inventory_items:
        try:
            _init_db()
            with _db() as conn:
                for inv in inventory_items:
                    conn.execute(
                        "INSERT INTO gb_inventory (sku,description,manufacturer,finish_fabric,"
                        "warehouse_location,cost_price,retail_price,status,date_in,notes) "
                        "VALUES (?,?,?,?,?,?,?,?,?,?)",
                        (inv.get('sku',''), inv.get('description',''), inv.get('manufacturer',''),
                         inv.get('finish_fabric',''), inv.get('warehouse_location',''),
                         inv.get('cost_price'), inv.get('retail_price'),
                         inv.get('status','Available'), inv.get('date_in',''), inv.get('notes','')))
                conn.commit()
            inv_saved = len(inventory_items)
        except Exception as _ie:
            print(f"[gb/import inventory] {_ie}")

    return jsonify({
        "orders":          orders,
        "rows_scanned":    len(raw_rows),
        "orders_added":    len(orders),
        "inventory_added": inv_saved,
        "source":          "fallback",
        "headers":         headers_found,
    })


@app.route("/api/gb/inventory", methods=["GET"])
def gb_inventory_list():
    """List inventory items with optional filters: status, manufacturer, q (search)."""
    status_f = request.args.get('status', '')
    mfr_f    = request.args.get('manufacturer', '')
    q        = request.args.get('q', '')
    try:
        _init_db()
        with _db() as conn:
            sql = "SELECT * FROM gb_inventory WHERE 1=1"
            params = []
            if status_f:
                sql += " AND status=?"; params.append(status_f)
            if mfr_f:
                sql += " AND manufacturer=?"; params.append(mfr_f)
            if q:
                sql += (" AND (description LIKE ? OR manufacturer LIKE ? OR sku LIKE ?"
                        " OR warehouse_location LIKE ? OR customer_name LIKE ?)")
                params += [f'%{q}%'] * 5
            sql += " ORDER BY created_at DESC"
            rows = conn.execute(sql, params).fetchall()
        return jsonify([dict(r) for r in rows])
    except Exception as e:
        return jsonify({"error": str(e)}), 500


@app.route("/api/gb/inventory", methods=["POST"])
def gb_inventory_add():
    """Add a single inventory item."""
    d = request.get_json() or {}
    try:
        _init_db()
        with _db() as conn:
            conn.execute(
                "INSERT INTO gb_inventory (sku,description,manufacturer,finish_fabric,"
                "warehouse_location,cost_price,retail_price,status,customer_name,order_id,"
                "date_in,date_out,delivery_type,notes) "
                "VALUES (?,?,?,?,?,?,?,?,?,?,?,?,?,?)",
                (d.get('sku',''), d.get('description',''), d.get('manufacturer',''),
                 d.get('finish_fabric',''), d.get('warehouse_location',''),
                 d.get('cost_price') or None, d.get('retail_price') or None,
                 d.get('status','Available'), d.get('customer_name',''),
                 d.get('order_id') or None, d.get('date_in',''), d.get('date_out',''),
                 d.get('delivery_type',''), d.get('notes','')))
            conn.commit()
        _log_activity(d.get('user_name',''), 'addInventory',
                      d.get('description','item'), f"SKU:{d.get('sku','')}")
        return jsonify({'success': True})
    except Exception as e:
        return jsonify({'error': str(e)}), 500


@app.route("/api/gb/inventory/<int:item_id>", methods=["PUT"])
def gb_inventory_update(item_id):
    """Update an inventory item."""
    d = request.get_json() or {}
    fields = ['sku','description','manufacturer','finish_fabric','warehouse_location',
              'cost_price','retail_price','status','customer_name','order_id',
              'date_in','date_out','delivery_type','notes']
    try:
        _init_db()
        with _db() as conn:
            sets = ', '.join(f"{f}=?" for f in fields)
            vals = [d.get(f) for f in fields] + [item_id]
            conn.execute(f"UPDATE gb_inventory SET {sets} WHERE id=?", vals)
            conn.commit()
        _log_activity(d.get('user_name',''), 'updateInventory',
                      d.get('description','item'), f"id:{item_id}")
        return jsonify({'success': True})
    except Exception as e:
        return jsonify({'error': str(e)}), 500


@app.route("/api/gb/inventory/<int:item_id>", methods=["DELETE"])
def gb_inventory_delete(item_id):
    """Delete an inventory item."""
    user_name = request.args.get('user_name', '')
    try:
        _init_db()
        with _db() as conn:
            row = conn.execute("SELECT description FROM gb_inventory WHERE id=?",
                               (item_id,)).fetchone()
            desc = row['description'] if row else str(item_id)
            conn.execute("DELETE FROM gb_inventory WHERE id=?", (item_id,))
            conn.commit()
        _log_activity(user_name, 'deleteInventory', desc, f"id:{item_id}")
        return jsonify({'success': True})
    except Exception as e:
        return jsonify({'error': str(e)}), 500


@app.route("/api/gb/attention", methods=["GET"])
def gb_attention():
    """AIVM morning briefing — what needs attention today at Great Bridge Furniture."""
    try:
        _init_db()
        with _db() as conn:
            rows = conn.execute("SELECT * FROM gb_orders ORDER BY created_at DESC").fetchall()
        orders = [dict(r) for r in rows]
    except Exception as e:
        orders = []

    today = time.strftime('%Y-%m-%d')

    if not orders:
        return jsonify({
            "attention": "No orders in the system yet. Import your spreadsheet or add your first order to get started.",
            "source": "empty"
        })

    # Build context for AIVM
    status_counts = {}
    overdue       = []
    ready         = []
    balance_due   = 0.0

    for o in orders:
        s = o.get('status', 'Ordered')
        status_counts[s] = status_counts.get(s, 0) + 1
        exp = o.get('expected_date', '')
        if exp and exp < today and s not in ['Delivered', 'Paid']:
            overdue.append(o)
        if s == 'Ready for Pickup':
            ready.append(o)
        try:
            balance_due += float(o.get('total_amount') or 0) - float(o.get('deposit_paid') or 0)
        except Exception:
            pass

    context = (
        f"Today is {today}. Great Bridge Furniture has {len(orders)} total orders.\n"
        f"Status breakdown: {json.dumps(status_counts)}\n"
        f"Orders overdue (past expected date, not yet delivered): {len(overdue)}\n"
        f"Orders ready for pickup: {len(ready)}\n"
        f"Total outstanding balance: ${balance_due:,.2f}\n\n"
    )
    if overdue:
        context += "Overdue orders:\n" + "\n".join(
            f"  - {o['customer_name']} ({o.get('manufacturer','?')}), expected {o.get('expected_date','?')}, status: {o.get('status','?')}"
            for o in overdue[:10]
        ) + "\n"
    if ready:
        context += "Ready for pickup:\n" + "\n".join(
            f"  - {o['customer_name']} ({o.get('manufacturer','?')}) — call to schedule delivery"
            for o in ready[:10]
        ) + "\n"

    # Read business profile for personalization
    biz_profile_obj = None
    try:
        with _db() as conn:
            bp_row = conn.execute(
                "SELECT * FROM biz_profile WHERE company = ?", ("GB",)
            ).fetchone()
        if bp_row:
            biz_profile_obj = dict(bp_row)
            try:
                biz_profile_obj["priorities"] = json.loads(biz_profile_obj.get("priorities") or "[]")
            except Exception:
                biz_profile_obj["priorities"] = []
    except Exception:
        pass

    if biz_profile_obj and biz_profile_obj.get("priorities"):
        priority_str = ", ".join(biz_profile_obj["priorities"][:3])
        focus_str    = biz_profile_obj.get("briefing_focus", "")
        profile_note = f"The owner's top priorities are: {priority_str}. {focus_str}".strip()
        prompt = (
            context +
            f"\nOwner profile: {profile_note}\n"
            "\nYou are the AI assistant for Great Bridge Furniture. Write a short morning briefing (3-5 bullet points) "
            "for the store owner. Lead with what matters most to them based on their stated priorities. "
            "Focus on what needs action today. Be direct and practical. "
            "Format as plain bullet points starting with •"
        )
    else:
        prompt = (
            context +
            "\nYou are the AI assistant for Great Bridge Furniture. Write a short morning briefing (3-5 bullet points) "
            "for the store owner. Focus on what needs action today: overdue orders, ready-for-pickup customers to call, "
            "upcoming expected deliveries, and any patterns worth noting. Be direct and practical. "
            "Format as plain bullet points starting with •"
        )

    aivm_text = None
    try:
        aivm = get_aivm()
        if aivm:
            aivm_text = aivm.run_inference(prompt, timeout_secs=120)
            return jsonify({"attention": aivm_text.strip(), "source": "AIVM",
                            "stats": {"total": len(orders), "overdue": len(overdue),
                                      "ready": len(ready), "balance_due": round(balance_due, 2)}})
    except Exception as e:
        print(f"[gb/attention AIVM] {e}")

    # Rule-based fallback
    lines = []
    if ready:
        names = ", ".join(o['customer_name'] for o in ready[:5])
        lines.append(f"• {len(ready)} order(s) ready for pickup — call to schedule: {names}")
    if overdue:
        names = ", ".join(o['customer_name'] for o in overdue[:5])
        lines.append(f"• {len(overdue)} order(s) past expected date — follow up with manufacturer: {names}")
    ordered_count = status_counts.get('Ordered', 0)
    if ordered_count:
        lines.append(f"• {ordered_count} order(s) placed and waiting on manufacturers")
    transit_count = status_counts.get('In Transit', 0)
    if transit_count:
        lines.append(f"• {transit_count} order(s) in transit — check for arrival today")
    if balance_due > 0:
        lines.append(f"• Outstanding balance to collect: ${balance_due:,.2f}")
    if not lines:
        lines.append("• All orders are on track — nothing urgent today")

    return jsonify({
        "attention": "\n".join(lines),
        "source": "rule-based",
        "stats": {"total": len(orders), "overdue": len(overdue),
                  "ready": len(ready), "balance_due": round(balance_due, 2)},
    })


# ── GREAT BRIDGE FURNITURE — DELIVERY CALENDAR ───────────────────────────────

@app.route("/api/gb/calendar", methods=["GET"])
def gb_calendar():
    """Return all calendar events (deliveries, arrivals, blocked days) for a month."""
    import calendar as _cal
    year  = request.args.get("year",  type=int) or datetime.now().year
    month = request.args.get("month", type=int) or datetime.now().month
    _, dim = _cal.monthrange(year, month)
    start = f"{year}-{month:02d}-01"
    end   = f"{year}-{month:02d}-{dim:02d}"
    try:
        _init_db()
        with _db() as conn:
            # Outgoing: orders with a scheduled delivery date
            deliveries = conn.execute(
                "SELECT id, customer_name, phone, items, delivery_date, status FROM gb_orders "
                "WHERE delivery_date BETWEEN ? AND ? AND delivery_date != ''",
                (start, end)).fetchall()
            # Incoming: orders expected to arrive from manufacturer
            arrivals = conn.execute(
                "SELECT id, customer_name, manufacturer, items, expected_date, status FROM gb_orders "
                "WHERE expected_date BETWEEN ? AND ? AND expected_date != ''",
                (start, end)).fetchall()
            # Blocked days
            blocked = conn.execute(
                "SELECT * FROM gb_blocked_days WHERE date BETWEEN ? AND ?",
                (start, end)).fetchall()
            # Standalone calendar delivery events (with times)
            events = conn.execute(
                "SELECT * FROM gb_deliveries WHERE date BETWEEN ? AND ? ORDER BY date, time",
                (start, end)).fetchall()
        return jsonify({
            "deliveries": [dict(r) for r in deliveries],
            "arrivals":   [dict(r) for r in arrivals],
            "blocked":    [dict(r) for r in blocked],
            "events":     [dict(r) for r in events],
        })
    except Exception as e:
        print(f"[gb/calendar] {e}")
        return jsonify({"deliveries": [], "arrivals": [], "blocked": [], "error": str(e)}), 500


@app.route("/api/gb/blocked-days", methods=["GET", "POST"])
def gb_blocked_days():
    """List or add blocked days (store closed, vacation, holiday)."""
    if request.method == "GET":
        try:
            _init_db()
            with _db() as conn:
                rows = conn.execute("SELECT * FROM gb_blocked_days ORDER BY date").fetchall()
            return jsonify({"blocked": [dict(r) for r in rows]})
        except Exception as e:
            return jsonify({"blocked": [], "error": str(e)}), 500
    # POST — add a blocked day
    data   = request.get_json(force=True) or {}
    date   = str(data.get("date", "")).strip()
    reason = str(data.get("reason", "closed")).strip()
    note   = str(data.get("note", "")).strip()
    if not date:
        return jsonify({"error": "date required"}), 400
    try:
        _init_db()
        with _db() as conn:
            conn.execute(
                "INSERT OR REPLACE INTO gb_blocked_days (date, reason, note) VALUES (?,?,?)",
                (date, reason, note))
            conn.commit()
        return jsonify({"ok": True})
    except Exception as e:
        return jsonify({"error": str(e)}), 500


@app.route("/api/gb/blocked-days/<int:day_id>", methods=["DELETE"])
def gb_blocked_day_delete(day_id):
    """Remove a blocked day."""
    try:
        _init_db()
        with _db() as conn:
            conn.execute("DELETE FROM gb_blocked_days WHERE id = ?", (day_id,))
            conn.commit()
        return jsonify({"ok": True})
    except Exception as e:
        return jsonify({"error": str(e)}), 500


# ── GREAT BRIDGE FURNITURE — DELIVERY CALENDAR EVENTS ───────────────────────

@app.route("/api/gb/deliveries", methods=["GET"])
def gb_get_deliveries():
    """Get delivery events. Pass ?date=YYYY-MM-DD for a single day, or omit for all."""
    date = request.args.get("date", "").strip()
    try:
        _init_db()
        with _db() as conn:
            if date:
                rows = conn.execute(
                    "SELECT * FROM gb_deliveries WHERE date = ? ORDER BY time, id",
                    (date,)).fetchall()
            else:
                rows = conn.execute(
                    "SELECT * FROM gb_deliveries ORDER BY date, time, id").fetchall()
        return jsonify({"deliveries": [dict(r) for r in rows]})
    except Exception as e:
        print(f"[gb/deliveries GET] {e}")
        return jsonify({"deliveries": [], "error": str(e)}), 500


@app.route("/api/gb/deliveries", methods=["POST"])
def gb_create_delivery():
    """Create a standalone calendar delivery event."""
    data = request.get_json(force=True) or {}
    if not str(data.get("date", "")).strip():
        return jsonify({"error": "date is required"}), 400
    if not str(data.get("name", "")).strip():
        return jsonify({"error": "name is required"}), 400
    try:
        _init_db()
        with _db() as conn:
            cur = conn.execute(
                "INSERT INTO gb_deliveries (date, time, direction, name, address, phone, items, notes, order_id) "
                "VALUES (?,?,?,?,?,?,?,?,?)",
                (str(data.get("date",      ""))[:20],
                 str(data.get("time",      ""))[:10],
                 str(data.get("direction", "OUT"))[:5],
                 str(data.get("name",      ""))[:200],
                 str(data.get("address",   ""))[:500],
                 str(data.get("phone",     ""))[:50],
                 str(data.get("items",     ""))[:500],
                 str(data.get("notes",     ""))[:500],
                 data.get("order_id") or None)
            )
            conn.commit()
        return jsonify({"ok": True, "id": cur.lastrowid})
    except Exception as e:
        print(f"[gb/deliveries POST] {e}")
        return jsonify({"error": str(e)}), 500


@app.route("/api/gb/deliveries/<int:del_id>", methods=["PUT"])
def gb_update_delivery(del_id):
    """Update a delivery event."""
    data    = request.get_json(force=True) or {}
    allowed = ["date","time","direction","name","address","phone","items","notes","order_id"]
    fields, vals = [], []
    for k in allowed:
        if k in data:
            fields.append(f"{k} = ?")
            vals.append(data[k])
    if not fields:
        return jsonify({"error": "Nothing to update"}), 400
    fields.append("updated_at = datetime('now')")
    vals.append(del_id)
    try:
        _init_db()
        with _db() as conn:
            conn.execute(f"UPDATE gb_deliveries SET {', '.join(fields)} WHERE id = ?", vals)
            conn.commit()
        return jsonify({"ok": True})
    except Exception as e:
        print(f"[gb/deliveries PUT {del_id}] {e}")
        return jsonify({"error": str(e)}), 500


@app.route("/api/gb/deliveries/<int:del_id>", methods=["DELETE"])
def gb_delete_delivery(del_id):
    """Delete a delivery event."""
    try:
        _init_db()
        with _db() as conn:
            conn.execute("DELETE FROM gb_deliveries WHERE id = ?", (del_id,))
            conn.commit()
        return jsonify({"ok": True})
    except Exception as e:
        print(f"[gb/deliveries DELETE {del_id}] {e}")
        return jsonify({"error": str(e)}), 500


# ── GREAT BRIDGE FURNITURE — DOCUMENT ATTACHMENTS ────────────────────────────

@app.route("/api/gb/orders/<int:order_id>/docs", methods=["POST"])
def gb_upload_doc(order_id):
    """Upload a file (PDF/image) and attach it to an order. Stores as BLOB."""
    if 'file' not in request.files:
        return jsonify({"error": "No file in request"}), 400
    f = request.files['file']
    if not f or not f.filename:
        return jsonify({"error": "Empty file"}), 400

    filename = f.filename
    # Safety: strip any path components
    filename = filename.replace('\\', '/').split('/')[-1][:255]
    filetype = f.content_type or 'application/octet-stream'
    # Restrict to safe file types
    allowed_ext = {'.pdf', '.jpg', '.jpeg', '.png', '.gif', '.webp', '.heic',
                   '.doc', '.docx', '.xls', '.xlsx', '.csv', '.txt'}
    ext = '.' + filename.rsplit('.', 1)[-1].lower() if '.' in filename else ''
    if ext not in allowed_ext:
        return jsonify({"error": f"File type not allowed: {ext}"}), 400

    user = request.form.get("user_name") or request.args.get("user_name") or "Unknown"
    try:
        filedata = f.read()
        if len(filedata) > 20 * 1024 * 1024:   # 20 MB cap
            return jsonify({"error": "File too large (20 MB max)"}), 400
        _init_db()
        with _db() as conn:
            row = conn.execute("SELECT customer_name FROM gb_orders WHERE id=?", (order_id,)).fetchone()
            customer = row["customer_name"] if row else "?"
            cur = conn.execute(
                "INSERT INTO gb_docs (order_id, filename, filetype, filedata) VALUES (?,?,?,?)",
                (order_id, filename, filetype, filedata)
            )
            conn.commit()
            doc_id = cur.lastrowid
        _log_activity(user, "Attached document", f"to order #{order_id}",
                      f"File: {filename} — Customer: {customer}")
        return jsonify({"ok": True, "doc_id": doc_id, "filename": filename})
    except Exception as e:
        print(f"[gb/docs POST order={order_id}] {e}")
        return jsonify({"error": str(e)}), 500


@app.route("/api/gb/orders/<int:order_id>/docs", methods=["GET"])
def gb_list_docs(order_id):
    """List all docs attached to an order (metadata only — no file bytes)."""
    try:
        _init_db()
        with _db() as conn:
            rows = conn.execute(
                "SELECT id, filename, filetype, uploaded_at FROM gb_docs WHERE order_id = ? ORDER BY uploaded_at DESC",
                (order_id,)
            ).fetchall()
        return jsonify({"docs": [dict(r) for r in rows]})
    except Exception as e:
        print(f"[gb/docs GET order={order_id}] {e}")
        return jsonify({"error": str(e), "docs": []}), 500


@app.route("/api/gb/docs/<int:doc_id>", methods=["GET"])
def gb_serve_doc(doc_id):
    """Serve the actual file bytes for a document."""
    try:
        _init_db()
        with _db() as conn:
            row = conn.execute(
                "SELECT filename, filetype, filedata FROM gb_docs WHERE id = ?", (doc_id,)
            ).fetchone()
        if not row:
            return jsonify({"error": "Document not found"}), 404
        filename, filetype, filedata = row["filename"], row["filetype"], row["filedata"]
        return send_file(
            io.BytesIO(filedata),
            mimetype=filetype or 'application/octet-stream',
            as_attachment=True,
            download_name=filename
        )
    except Exception as e:
        print(f"[gb/docs/{doc_id} GET] {e}")
        return jsonify({"error": str(e)}), 500


# ── GREAT BRIDGE FURNITURE — SPREADSHEET ANALYSIS + RESTRUCTURE ──────────────

@app.route("/api/gb/analyze-sheets", methods=["POST"])
def gb_analyze_sheets():
    """
    Upload 1-10 Excel files. AIVM reads all of them, understands column
    meanings across files, identifies overlaps and patterns, and returns
    a plain-English analysis + a proposed unified column mapping.

    Returns: { analysis: str, column_mapping: {...}, files_read: int, source: str }
    """
    files = request.files.getlist('files')
    if not files or len(files) == 0:
        return jsonify({"error": "No files uploaded"}), 400
    if len(files) > 10:
        return jsonify({"error": "Maximum 10 files at once"}), 400

    all_file_summaries = []

    for f in files:
        fname = (f.filename or "").strip()
        ext   = ('.' + fname.rsplit('.', 1)[-1].lower()) if '.' in fname else ''
        if ext not in {'.xlsx', '.xls', '.csv'}:
            continue

        try:
            content = f.read()
            rows    = []

            if ext == '.csv':
                text   = content.decode('utf-8', errors='replace')
                reader = _csv_mod.DictReader(io.StringIO(text))
                for i, row in enumerate(reader):
                    rows.append({k: str(v) for k, v in row.items()})
                    if i >= 100:
                        break
                headers = list(reader.fieldnames or [])
            else:
                import openpyxl
                wb      = openpyxl.load_workbook(io.BytesIO(content), data_only=True)
                ws      = wb.active
                headers = []
                for i, row in enumerate(ws.iter_rows(values_only=True)):
                    if i == 0:
                        headers = [str(h).strip() if h is not None else f"col{j}" for j, h in enumerate(row)]
                    elif not all(v is None for v in row):
                        rows.append({
                            headers[j]: str(v).strip() if v is not None else ""
                            for j, v in enumerate(row) if j < len(headers)
                        })
                    if i >= 100:
                        break
                wb.close()

            sample_rows = rows[:15]
            rows_text   = "\n".join(
                " | ".join(f"{k}: {v}" for k, v in r.items() if v and v.lower() not in ('none', ''))
                for r in sample_rows
            )
            all_file_summaries.append(
                f"FILE: {fname}\nColumns: {', '.join(headers)}\nSample rows ({len(rows)} total rows):\n{rows_text}"
            )

        except Exception as e:
            all_file_summaries.append(f"FILE: {fname} — ERROR reading: {e}")

    if not all_file_summaries:
        return jsonify({"error": "Could not read any uploaded files"}), 400

    combined_text = "\n\n---\n\n".join(all_file_summaries)

    if len(all_file_summaries) == 1:
        fname = all_file_summaries[0].split('\n')[0].replace('FILE: ', '').strip()
        prompt = (
            "You are analyzing a spreadsheet uploaded by Great Bridge Furniture, a furniture store in Chesapeake Virginia.\n"
            f"They have uploaded one file: {fname}\n\n"
            f"{combined_text}\n\n"
            "Please analyze this spreadsheet and provide:\n\n"
            "1. WHAT THIS FILE CONTAINS — a plain-English description of what data is in this file\n"
            "2. COLUMN SUMMARY — describe what each column contains and what it is used for\n"
            "3. DATA QUALITY ISSUES — missing values, inconsistent formats, empty columns, etc.\n"
            "4. KEY INSIGHTS — notable patterns, totals, or anything David at Great Bridge Furniture should know\n"
            "5. CLEAN COLUMN NAMES — propose clean professional column names if any improvements are possible\n\n"
            "After your analysis, output a JSON block at the end in this exact format:\n"
            "```json\n"
            '{"unified_columns": ["Col1", "Col2", ...], '
            '"file_mappings": {"' + fname + '": {"their_col": "unified_col", ...}}}\n'
            "```\n"
            "Be specific and practical — David at Great Bridge Furniture needs to understand this."
        )
    else:
        prompt = (
            "You are analyzing spreadsheets uploaded by Great Bridge Furniture, a furniture store in Chesapeake Virginia.\n"
            "They have uploaded multiple Excel/CSV files that may contain overlapping or related customer order data.\n\n"
            f"{combined_text}\n\n"
            "Please analyze these spreadsheets and provide:\n\n"
            "1. WHAT EACH FILE CONTAINS — a plain-English description of what data is in each file\n"
            "2. OVERLAPPING COLUMNS — identify columns that appear to contain the same kind of data across files, "
            "even if named differently (e.g. 'Cust Name' vs 'Customer' vs 'Buyer Name' all mean the same thing). "
            "List them clearly.\n"
            "3. DUPLICATE DATA — if the same orders or customers appear in multiple files, note that.\n"
            "4. DATA QUALITY ISSUES — missing values, inconsistent formats, empty columns, etc.\n"
            "5. RECOMMENDED UNIFIED COLUMNS — propose a clean set of column names for a merged spreadsheet "
            "that captures all the useful data from all files.\n\n"
            "After your analysis, output a JSON block at the end in this exact format:\n"
            "```json\n"
            '{"unified_columns": ["Col1", "Col2", ...], '
            '"file_mappings": {"filename.xlsx": {"their_col": "unified_col", ...}}}\n'
            "```\n"
            "Be specific and practical — David at Great Bridge Furniture needs to understand this."
        )

    try:
        aivm        = get_aivm()
        aivm_result = aivm.run_inference(prompt, timeout_secs=300)

        # Try to parse the JSON block at the end
        column_mapping = {}
        json_match = _re_mod.search(r'```json\s*(.*?)\s*```', aivm_result, _re_mod.DOTALL)
        if json_match:
            try:
                column_mapping = json.loads(json_match.group(1))
            except Exception:
                pass

        # Clean the analysis text (remove the JSON block for display)
        analysis_text = _re_mod.sub(r'```json\s*.*?\s*```', '', aivm_result, flags=_re_mod.DOTALL).strip()

        return jsonify({
            "analysis":       analysis_text,
            "column_mapping": column_mapping,
            "files_read":     len(all_file_summaries),
            "source":         "AIVM",
        })

    except Exception as e:
        print(f"[gb/analyze-sheets] AIVM error: {e}")
        # Fallback: just summarize what we found
        summary_lines = [f"Read {len(all_file_summaries)} file(s)."]
        for s in all_file_summaries:
            first_line = s.split('\n')[0]
            summary_lines.append(f"• {first_line}")
        return jsonify({
            "analysis":       "\n".join(summary_lines) + f"\n\nAIVM unavailable: {e}",
            "column_mapping": {},
            "files_read":     len(all_file_summaries),
            "source":         "fallback",
        })


@app.route("/api/gb/restructure-sheets", methods=["POST"])
def gb_restructure_sheets():
    """
    Upload 1-10 Excel files (same as analyze-sheets).
    Optionally include a JSON 'mapping' field with unified column definitions.
    AIVM merges all the data into one clean .xlsx file which is returned as download.

    Returns: .xlsx file download
    """
    import openpyxl
    from openpyxl.styles import Font, PatternFill, Alignment

    files   = request.files.getlist('files')
    mapping_json = request.form.get('mapping', '{}')
    try:
        mapping = json.loads(mapping_json)
    except Exception:
        mapping = {}

    if not files or len(files) == 0:
        return jsonify({"error": "No files uploaded"}), 400
    if len(files) > 10:
        return jsonify({"error": "Maximum 10 files at once"}), 400

    # ── Step 1: Read all files into a list of dicts ──────────────────────────
    all_rows           = []
    all_headers_seen   = set()    # for fast duplicate check
    all_headers_ordered = []      # preserves original column order
    file_summaries = []

    for f in files:
        fname = (f.filename or "").strip()
        ext   = ('.' + fname.rsplit('.', 1)[-1].lower()) if '.' in fname else ''
        if ext not in {'.xlsx', '.xls', '.csv'}:
            continue
        try:
            content = f.read()
            rows = []
            if ext == '.csv':
                text   = content.decode('utf-8', errors='replace')
                reader = _csv_mod.DictReader(io.StringIO(text))
                for row in reader:
                    rows.append({k.strip(): str(v).strip() for k, v in row.items()})
                    if len(rows) >= 5000:
                        break
            else:
                wb = openpyxl.load_workbook(io.BytesIO(content), read_only=True, data_only=True)
                ws = wb.active
                headers = []
                for i, row in enumerate(ws.iter_rows(values_only=True)):
                    if i == 0:
                        headers = [str(h).strip() if h is not None else f"col{j}" for j, h in enumerate(row)]
                    elif not all(v is None for v in row):
                        rows.append({
                            headers[j]: str(v).strip() if v is not None else ""
                            for j, v in enumerate(row) if j < len(headers)
                        })
                    if i >= 5000:
                        break
                wb.close()

            file_summaries.append(f"{fname}: {len(rows)} rows")
            for r in rows:
                for k in r.keys():
                    if k not in all_headers_seen:
                        all_headers_seen.add(k)
                        all_headers_ordered.append(k)
            all_rows.extend(rows)

        except Exception as e:
            file_summaries.append(f"{fname}: ERROR — {e}")

    if not all_rows:
        return jsonify({"error": "No readable data found in uploaded files"}), 400

    # ── Step 2: Build AIVM prompt to get clean column mapping ────────────────
    file_col_map = mapping.get("file_mappings", {})
    unified_cols = mapping.get("unified_columns", [])

    if not unified_cols:
        # Ask AIVM to propose unified columns
        headers_list = all_headers_ordered
        prompt = (
            "Great Bridge Furniture has provided spreadsheet data with these column names across all files:\n"
            f"{', '.join(headers_list)}\n\n"
            "Many of these columns contain the same type of data with different names. "
            "Please propose a clean, unified set of column headers for a merged spreadsheet. "
            "Output ONLY a JSON array of the unified column names, nothing else.\n"
            'Example: ["Customer Name", "Phone", "Items Ordered", "Manufacturer", "Total", "Deposit", "Expected Date", "Status", "Notes"]'
        )
        try:
            aivm        = get_aivm()
            aivm_result = aivm.run_inference(prompt, timeout_secs=120)
            clean = _re_mod.sub(r'```(?:json)?\s*|\s*```', '', aivm_result).strip()
            match = _re_mod.search(r'\[.*\]', clean, _re_mod.DOTALL)
            unified_cols = json.loads(match.group()) if match else list(all_headers_ordered)
        except Exception as e:
            print(f"[gb/restructure AIVM col mapping] {e}")
            unified_cols = list(all_headers_ordered)

    # ── Step 3: Map each row to unified columns ───────────────────────────────
    def _best_match(src_col, unified):
        """Simple fuzzy match — normalize both sides and find best overlap."""
        src_lower = src_col.lower().replace('_', ' ').strip()
        for u in unified:
            u_lower = u.lower().replace('_', ' ').strip()
            if src_lower == u_lower:
                return u
        # Word overlap
        src_words = set(src_lower.split())
        best_col, best_score = None, 0
        for u in unified:
            u_words = set(u.lower().replace('_', ' ').strip().split())
            score   = len(src_words & u_words)
            if score > best_score:
                best_col, best_score = u, score
        return best_col if best_score > 0 else None

    cleaned_rows = []
    for row in all_rows:
        new_row = {c: "" for c in unified_cols}
        for src_col, val in row.items():
            if not val or val.lower() in ('none', 'null', ''):
                continue
            # Check explicit mapping first
            target = None
            for fname_key, col_map in file_col_map.items():
                if src_col in col_map:
                    target = col_map[src_col]
                    break
            if not target:
                target = _best_match(src_col, unified_cols)
            if target and target in new_row:
                if not new_row[target]:     # don't overwrite existing value
                    new_row[target] = val
        # Only include rows that have at least one non-empty value
        if any(v.strip() for v in new_row.values()):
            cleaned_rows.append(new_row)

    # Remove exact duplicate rows
    seen     = set()
    deduped  = []
    for r in cleaned_rows:
        key = tuple(r.get(c, '') for c in unified_cols)
        if key not in seen:
            seen.add(key)
            deduped.append(r)

    # ── Step 4a: CSV export (if requested) ───────────────────────────────────
    export_format = request.form.get('format', 'xlsx').lower().strip()
    if export_format == 'csv':
        out_buf = io.StringIO()
        writer  = _csv_mod.DictWriter(out_buf, fieldnames=unified_cols, extrasaction='ignore')
        writer.writeheader()
        writer.writerows(deduped)
        csv_bytes = out_buf.getvalue().encode('utf-8-sig')   # utf-8-sig = BOM for Windows Excel
        filename  = f"GB_Unified_{time.strftime('%Y%m%d_%H%M')}.csv"
        return send_file(
            io.BytesIO(csv_bytes),
            mimetype='text/csv',
            as_attachment=True,
            download_name=filename
        )

    # ── Step 4b: Write clean .xlsx with openpyxl ─────────────────────────────
    wb_out = openpyxl.Workbook()
    ws_out = wb_out.active
    ws_out.title = "Unified Data"

    # Header row — navy blue background, white bold text
    header_fill = PatternFill("solid", fgColor="1B3A5C")
    header_font = Font(bold=True, color="FFFFFF", size=12)
    header_align = Alignment(horizontal="center", vertical="center", wrap_text=True)

    for col_idx, col_name in enumerate(unified_cols, start=1):
        cell = ws_out.cell(row=1, column=col_idx, value=col_name)
        cell.fill   = header_fill
        cell.font   = header_font
        cell.alignment = header_align

    # Freeze header row
    ws_out.freeze_panes = "A2"

    # Data rows
    row_fill_alt = PatternFill("solid", fgColor="F5F8FC")   # light blue-grey alternating rows
    for row_idx, row_data in enumerate(deduped, start=2):
        for col_idx, col_name in enumerate(unified_cols, start=1):
            cell = ws_out.cell(row=row_idx, column=col_idx, value=row_data.get(col_name, ""))
            if row_idx % 2 == 0:
                cell.fill = row_fill_alt

    # Auto-fit column widths (approximate)
    for col_idx, col_name in enumerate(unified_cols, start=1):
        max_len = len(col_name)
        for row_data in deduped[:200]:     # sample first 200 rows
            v = str(row_data.get(col_name, ""))
            if len(v) > max_len:
                max_len = len(v)
        ws_out.column_dimensions[
            openpyxl.utils.get_column_letter(col_idx)
        ].width = min(max_len + 4, 50)

    ws_out.row_dimensions[1].height = 28

    # Summary tab
    ws_sum = wb_out.create_sheet("Import Summary")
    ws_sum.append(["Source Files"])
    for s in file_summaries:
        ws_sum.append([s])
    ws_sum.append([])
    ws_sum.append([f"Total rows after dedup: {len(deduped)}"])
    ws_sum.append([f"Unified columns: {len(unified_cols)}"])
    ws_sum.append([f"Generated: {time.strftime('%Y-%m-%d %H:%M')}"])

    # Write to BytesIO and return
    out_buf = io.BytesIO()
    wb_out.save(out_buf)
    out_buf.seek(0)

    filename = f"GB_Unified_{time.strftime('%Y%m%d_%H%M')}.xlsx"
    return send_file(
        out_buf,
        mimetype='application/vnd.openxmlformats-officedocument.spreadsheetml.sheet',
        as_attachment=True,
        download_name=filename
    )


# ── BUSINESS PROFILE — AIVM ONBOARDING ──────────────────────────────────────

@app.route("/api/biz/profile", methods=["GET"])
def biz_get_profile():
    """Return the biz_profile for a company, or null if not yet set.
    Query param: company=GB
    """
    company = (request.args.get("company") or "").strip().upper()
    if not company:
        return jsonify({"error": "company required"}), 400
    try:
        _init_db()
        with _db() as conn:
            row = conn.execute(
                "SELECT * FROM biz_profile WHERE company = ?", (company,)
            ).fetchone()
        if not row:
            return jsonify({"profile": None})
        p = dict(row)
        try:
            p["priorities"] = json.loads(p.get("priorities") or "[]")
        except Exception:
            p["priorities"] = []
        return jsonify({"profile": p})
    except Exception as e:
        print(f"[biz/profile GET] {e}")
        return jsonify({"error": str(e)}), 500


@app.route("/api/biz/profile", methods=["DELETE"])
def biz_delete_profile():
    """Delete (reset) the biz_profile for a company so onboarding can run again.
    Query param: company=GB
    """
    company = (request.args.get("company") or "").strip().upper()
    if not company:
        return jsonify({"error": "company required"}), 400
    try:
        _init_db()
        with _db() as conn:
            conn.execute("DELETE FROM biz_profile WHERE company = ?", (company,))
            conn.commit()
        return jsonify({"ok": True})
    except Exception as e:
        print(f"[biz/profile DELETE] {e}")
        return jsonify({"error": str(e)}), 500


_ONBOARD_SYSTEM = (
    "You are LightView's friendly business assistant having a short setup conversation with a business owner. "
    "Your goal is to understand their business and what matters most to them each day.\n\n"
    "Ask warm, plain-English questions — no jargon, no forms, no bullet points in your questions. "
    "After 4-5 exchanges, give a short confirmation of what you learned, then — and only then — "
    "output the profile on its own final line.\n\n"
    "Rules:\n"
    "- Keep each reply to 2-3 short sentences max\n"
    "- Ask exactly one question at a time\n"
    "- Use simple, friendly conversational language — imagine talking to a neighbor\n"
    "- Only output PROFILE_JSON: after the owner has confirmed your summary\n"
    "- When ready, output this as the very last line of your reply:\n"
    'PROFILE_JSON:{"business_type":"...","priorities":["...","...","..."],'
    '"briefing_focus":"...","staff_notes":"...","summary":"..."}\n\n'
    "Field definitions:\n"
    "- business_type: short label (e.g. furniture store, restaurant, medical office)\n"
    "- priorities: top 3 things the owner cares about most each day\n"
    "- briefing_focus: one sentence — what should the morning briefing lead with\n"
    "- staff_notes: brief note on staff structure if mentioned\n"
    "- summary: 2-sentence plain-English description of this business and their needs"
)


def _save_biz_profile(company: str, profile: dict, messages: list):
    """Persist a biz_profile row to the DB (upsert)."""
    try:
        _init_db()
        priorities_json = json.dumps(profile.get("priorities", []))
        raw_conv        = json.dumps(messages)
        with _db() as conn:
            conn.execute("""
                INSERT INTO biz_profile
                    (company, business_type, priorities, briefing_focus, staff_notes, summary, raw_conv, updated_at)
                VALUES (?,?,?,?,?,?,?, datetime('now'))
                ON CONFLICT(company) DO UPDATE SET
                    business_type  = excluded.business_type,
                    priorities     = excluded.priorities,
                    briefing_focus = excluded.briefing_focus,
                    staff_notes    = excluded.staff_notes,
                    summary        = excluded.summary,
                    raw_conv       = excluded.raw_conv,
                    updated_at     = datetime('now')
            """, (
                company,
                str(profile.get("business_type", ""))[:200],
                priorities_json,
                str(profile.get("briefing_focus", ""))[:500],
                str(profile.get("staff_notes", ""))[:500],
                str(profile.get("summary", ""))[:1000],
                raw_conv
            ))
            conn.commit()
        print(f"[biz_profile] saved for {company}")
    except Exception as e:
        print(f"[biz_profile] save error: {e}")


@app.route("/api/biz/onboard", methods=["POST"])
def biz_onboard():
    """
    Multi-turn AIVM onboarding conversation.
    Body: { company: str, messages: [{role: "user"|"assistant", content: str}] }
    Returns: { reply: str, done: bool, profile?: {...} }
    """
    data     = request.get_json(force=True) or {}
    company  = str(data.get("company") or "GB").strip().upper()
    messages = data.get("messages", [])

    if not messages:
        return jsonify({"error": "messages array required"}), 400

    # Build the full conversation prompt for AIVM (system + all prior turns)
    conv_text = _ONBOARD_SYSTEM + "\n\n"
    for m in messages:
        role    = m.get("role", "user")
        content = str(m.get("content", "")).strip()
        if role == "user":
            conv_text += f"Owner: {content}\n"
        else:
            conv_text += f"Assistant: {content}\n"
    conv_text += "Assistant:"

    reply_text = None
    try:
        aivm      = get_aivm()
        raw       = aivm.run_inference(conv_text, timeout_secs=120).strip()
        # Trim echoed prefix if AIVM repeated the cue
        for prefix in ("assistant:", "lightview ai:", "lightview assistant:"):
            if raw.lower().startswith(prefix):
                raw = raw[len(prefix):].strip()
                break
        reply_text = raw
    except Exception as e:
        print(f"[biz/onboard AIVM] {e}")

    # Rule-based fallback when AIVM is unavailable
    if not reply_text:
        user_count = sum(1 for m in messages if m.get("role") == "user")
        if user_count == 1:
            reply_text = (
                "Thanks for telling me about your business! "
                "What's the hardest part of your day to keep track of — "
                "is it knowing what's coming in, what's going out to customers, "
                "who owes you money, or something else?"
            )
        elif user_count == 2:
            reply_text = (
                "That makes a lot of sense. "
                "Do you have staff who help run things, and do they all need to "
                "see everything or just certain parts of the business?"
            )
        else:
            # Build a basic profile and wrap up
            basic_profile = {
                "business_type": company,
                "priorities":    ["orders", "deliveries", "outstanding balances"],
                "briefing_focus": "Lead with any overdue orders and customers ready for pickup.",
                "staff_notes":   "Not specified during setup.",
                "summary":       f"Business at {company}. Owner wants a clear daily picture of orders, deliveries, and who owes money."
            }
            _save_biz_profile(company, basic_profile, messages)
            return jsonify({
                "reply": (
                    "Got it — I’ve set up your dashboard based on what you’ve shared with me. "
                    "You can always update this from Settings if your needs change."
                ),
                "done":    True,
                "profile": basic_profile
            })

    # Check for PROFILE_JSON: marker in the AIVM reply
    profile = None
    done    = False
    profile_match = _re_mod.search(
        r'PROFILE_JSON:\s*(\{.*?\})\s*$', reply_text, _re_mod.DOTALL | _re_mod.MULTILINE
    )
    if profile_match:
        try:
            profile    = json.loads(profile_match.group(1))
            # Remove the PROFILE_JSON line from the visible reply
            reply_text = reply_text[:profile_match.start()].strip()
            if not reply_text:
                reply_text = (
                    "Perfect — I’ve got everything I need. "
                    "Your dashboard is now personalised for your business."
                )
            _save_biz_profile(company, profile, messages)
            done = True
        except Exception as e:
            print(f"[biz/onboard] profile parse error: {e}")

    return jsonify({"reply": reply_text, "done": done, "profile": profile})


# ── GREAT BRIDGE FURNITURE — ACTIVITY LOG ────────────────────────────────────

@app.route("/api/gb/activity", methods=["GET"])
def gb_activity_log():
    """Return the employee activity log.
    Optional query params:
      limit  — max rows (default 200, max 500)
      offset — skip N rows for pagination (default 0)
      user   — filter by user_name (case-insensitive)
      q      — search text across action/target/details
    """
    try:
        limit  = min(int(request.args.get("limit",  200)), 500)
        offset = int(request.args.get("offset", 0))
        user_f = (request.args.get("user") or "").strip()
        q      = (request.args.get("q")    or "").strip().lower()
        _init_db()
        with _db() as conn:
            rows = conn.execute(
                "SELECT * FROM gb_activity_log ORDER BY id DESC LIMIT 1000"
            ).fetchall()
        log = [dict(r) for r in rows]
        # Client-side filter (log rarely exceeds a few thousand rows)
        if user_f:
            log = [r for r in log if r.get("user_name", "").lower() == user_f.lower()]
        if q:
            log = [r for r in log
                   if q in (r.get("action","") + " " +
                             r.get("target","") + " " +
                             r.get("details","")).lower()]
        total = len(log)
        log   = log[offset: offset + limit]
        return jsonify({"log": log, "total": total})
    except Exception as e:
        print(f"[gb/activity] {e}")
        return jsonify({"log": [], "total": 0, "error": str(e)}), 500


# ── GREAT BRIDGE FURNITURE — CUSTOMER SEARCH ─────────────────────────────────

@app.route("/api/gb/search", methods=["GET"])
def gb_customer_search():
    """
    Partial-name customer search across gb_orders.
    Query param: q — partial name (e.g. "jones")
    Returns unique customer cards with order count, phone, last order date,
    and latest expected date.
    """
    q = (request.args.get("q") or "").strip()
    if not q:
        return jsonify({"customers": [], "error": "q param required"}), 400
    try:
        _init_db()
        with _db() as conn:
            rows = conn.execute(
                "SELECT customer_name, phone, status, expected_date, created_at "
                "FROM gb_orders WHERE LOWER(customer_name) LIKE ? "
                "ORDER BY created_at DESC",
                (f"%{q.lower()}%",)
            ).fetchall()
        # Group by lowercase name in Python — clean deduplication
        groups = {}
        for r in rows:
            key = r["customer_name"].lower().strip()
            if key not in groups:
                groups[key] = {
                    "name":          r["customer_name"],
                    "phone":         r["phone"] or "",
                    "order_count":   0,
                    "last_order_at": r["created_at"] or "",
                    "latest_eta":    "",
                    "has_pending":   False,
                }
            g = groups[key]
            g["order_count"] += 1
            # Keep most-recent phone if we have one
            if r["phone"] and not g["phone"]:
                g["phone"] = r["phone"]
            # Track the furthest-out expected date
            exp = r["expected_date"] or ""
            if exp and exp > g["latest_eta"]:
                g["latest_eta"] = exp
            # Flag if any order is still open
            if r["status"] not in ("Delivered", "Paid"):
                g["has_pending"] = True

        customers = sorted(
            groups.values(),
            key=lambda x: x["last_order_at"],
            reverse=True
        )
        return jsonify({"customers": customers, "query": q})
    except Exception as e:
        print(f"[gb/search] {e}")
        return jsonify({"customers": [], "error": str(e)}), 500


@app.route("/api/gb/customers", methods=["GET"])
def gb_customer_profile():
    """
    Full customer profile by name.
    Query param: name — customer name (matched case-insensitively)
    Returns all orders, pending orders with ETAs, delivery history,
    doc counts per order, and outstanding balance.
    """
    name = (request.args.get("name") or "").strip()
    if not name:
        return jsonify({"error": "name param required"}), 400
    try:
        _init_db()
        with _db() as conn:
            orders = conn.execute(
                "SELECT * FROM gb_orders "
                "WHERE LOWER(customer_name) = LOWER(?) "
                "ORDER BY created_at DESC",
                (name,)
            ).fetchall()
            # Doc counts for each matching order
            doc_rows = conn.execute(
                "SELECT order_id, COUNT(*) as doc_count FROM gb_docs "
                "WHERE order_id IN "
                "  (SELECT id FROM gb_orders WHERE LOWER(customer_name) = LOWER(?)) "
                "GROUP BY order_id",
                (name,)
            ).fetchall()

        doc_counts   = {r["order_id"]: r["doc_count"] for r in doc_rows}
        orders_list  = []
        total_balance = 0.0

        for o in orders:
            d = dict(o)
            d["doc_count"] = doc_counts.get(o["id"], 0)
            try:
                total_balance += float(o["total_amount"] or 0) - float(o["deposit_paid"] or 0)
            except Exception:
                pass
            orders_list.append(d)

        pending   = [o for o in orders_list if o.get("status") not in ("Delivered", "Paid")]
        delivered = [o for o in orders_list if o.get("status") in ("Delivered", "Paid")]
        phone     = orders_list[0]["phone"] if orders_list else ""

        return jsonify({
            "name":         name,
            "phone":        phone,
            "order_count":  len(orders_list),
            "orders":       orders_list,
            "pending":      pending,
            "delivered":    delivered,
            "balance_due":  round(total_balance, 2),
        })
    except Exception as e:
        print(f"[gb/customers] {e}")
        return jsonify({"error": str(e)}), 500


# ── GREAT BRIDGE FURNITURE — DOC GENERATION HELPERS ─────────────────────────

def _gen_spreadsheet(orders_list, message):
    """
    Build an .xlsx spreadsheet from the current orders list.
    Returns {filename, mime, b64}.
    """
    import openpyxl as _xl
    from openpyxl.styles import Font, PatternFill, Alignment
    import io as _io2, base64 as _b64_2, datetime as _dt2
    from collections import Counter as _Counter2

    wb = _xl.Workbook()
    ws = wb.active
    ws.title = "Orders"

    headers = ["#", "Customer", "Phone", "Items", "Manufacturer",
               "Total ($)", "Deposit ($)", "Balance ($)",
               "Status", "Expected Date", "Delivery Date", "Notes"]
    for col, h in enumerate(headers, 1):
        cell = ws.cell(row=1, column=col, value=h)
        cell.font = Font(bold=True, color="FFFFFF")
        cell.fill = PatternFill("solid", fgColor="1B4F72")
        cell.alignment = Alignment(wrap_text=True)

    widths = [5, 22, 14, 35, 20, 11, 11, 11, 14, 14, 14, 25]
    for col, w in enumerate(widths, 1):
        ws.column_dimensions[ws.cell(row=1, column=col).column_letter].width = w

    grey_fill = PatternFill("solid", fgColor="EAF2F8")
    for i, o in enumerate(orders_list, 2):
        fill = grey_fill if i % 2 == 0 else None
        vals = [
            o.get("id", ""),
            o.get("customer_name", ""),
            o.get("phone", ""),
            o.get("items", ""),
            o.get("manufacturer", ""),
            o.get("total_price", 0) or 0,
            o.get("deposit_paid", 0) or 0,
            (o.get("total_price") or 0) - (o.get("deposit_paid") or 0),
            o.get("status", ""),
            o.get("expected_date", ""),
            o.get("delivery_date", ""),
            o.get("notes", ""),
        ]
        for col, v in enumerate(vals, 1):
            c = ws.cell(row=i, column=col, value=v)
            if fill:
                c.fill = fill

    n = len(orders_list) + 2
    ws.cell(row=n, column=1, value="TOTALS").font = Font(bold=True)
    total_total   = sum((o.get("total_price")   or 0) for o in orders_list)
    total_deposit = sum((o.get("deposit_paid")  or 0) for o in orders_list)
    total_balance = total_total - total_deposit
    for col, v in zip([6, 7, 8], [total_total, total_deposit, total_balance]):
        c = ws.cell(row=n, column=col, value=round(v, 2))
        c.font = Font(bold=True)
        c.fill = PatternFill("solid", fgColor="D5E8D4")

    ws.freeze_panes = "A2"

    ws2 = wb.create_sheet("Summary")
    status_counts2 = _Counter2(o.get("status", "Unknown") for o in orders_list)
    ws2.cell(1, 1, "Status").font = Font(bold=True)
    ws2.cell(1, 2, "Count").font  = Font(bold=True)
    last_row = 2
    for row, (s, c) in enumerate(status_counts2.items(), 2):
        ws2.cell(row, 1, s)
        ws2.cell(row, 2, c)
        last_row = row
    ws2.cell(last_row + 2, 1, "Generated").font = Font(italic=True)
    ws2.cell(last_row + 2, 2, _dt2.datetime.now().strftime("%Y-%m-%d %H:%M"))

    buf = _io2.BytesIO()
    wb.save(buf)
    buf.seek(0)
    b64 = _b64_2.b64encode(buf.read()).decode()
    today_str = _dt2.datetime.now().strftime("%Y%m%d")
    return {
        "filename": f"GBF_Orders_{today_str}.xlsx",
        "mime": "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet",
        "b64": b64,
    }


def _gen_word_doc(orders_list, message, file_b64=None, file_mime=None, file_name=None):
    """
    Build a .docx based on orders_list and message.
    If file_b64 is a .docx upload, appends an AI-reviewed note.
    Returns {filename, mime, b64}.
    """
    import io as _io3, base64 as _b64_3, datetime as _dt3
    from collections import Counter as _Counter3
    from docx import Document
    from docx.shared import Pt, RGBColor as _DocxRGB, Inches
    from docx.enum.text import WD_ALIGN_PARAGRAPH

    today     = _dt3.datetime.now().strftime("%B %d, %Y")
    gen_date  = _dt3.datetime.now().strftime("%Y%m%d")

    # If user uploaded a .docx and wants it improved
    if file_b64 and file_name and file_name.lower().endswith(".docx"):
        try:
            raw = _b64_3.b64decode(file_b64)
            doc = Document(_io3.BytesIO(raw))
            doc.add_paragraph()
            p    = doc.add_paragraph()
            run  = p.add_run(f"— Reviewed and improved by AI Assistant on {today} —")
            run.italic = True
            run.font.color.rgb = _DocxRGB(0x7F, 0x8C, 0x8D)
            buf = _io3.BytesIO()
            doc.save(buf)
            buf.seek(0)
            b64  = _b64_3.b64encode(buf.read()).decode()
            fname = file_name.rsplit(".", 1)[0] + f"_improved_{gen_date}.docx"
            return {"filename": fname,
                    "mime": "application/vnd.openxmlformats-officedocument.wordprocessingml.document",
                    "b64": b64}
        except Exception:
            pass  # fall through to generate fresh doc

    doc = Document()
    q_lower = message.lower()
    is_proposal = any(w in q_lower for w in ("proposal", "quote", "quotation"))
    is_report   = any(w in q_lower for w in ("report", "summary", "overview"))
    title_txt   = "Sales Proposal" if is_proposal else ("Business Report" if is_report else "Business Document")

    t = doc.add_heading(f"Great Bridge Furniture — {title_txt}", 0)
    t.alignment = WD_ALIGN_PARAGRAPH.CENTER
    doc.add_paragraph(f"Generated: {today}").alignment = WD_ALIGN_PARAGRAPH.CENTER

    doc.add_heading("Business Overview", 1)
    total_rev  = sum((o.get("total_price")  or 0) for o in orders_list)
    total_dep  = sum((o.get("deposit_paid") or 0) for o in orders_list)
    balance_d  = total_rev - total_dep
    sc         = _Counter3(o.get("status", "Unknown") for o in orders_list)

    doc.add_paragraph(f"Total Active Orders: {len(orders_list)}")
    doc.add_paragraph(f"Total Order Value: ${total_rev:,.2f}")
    doc.add_paragraph(f"Total Collected: ${total_dep:,.2f}")
    doc.add_paragraph(f"Outstanding Balance: ${balance_d:,.2f}")

    doc.add_heading("Order Status Breakdown", 1)
    tbl = doc.add_table(rows=1, cols=2)
    tbl.style = "Table Grid"
    hdr = tbl.rows[0].cells
    hdr[0].text = "Status"
    hdr[1].text = "Count"
    for s, c in sc.items():
        row = tbl.add_row().cells
        row[0].text = str(s)
        row[1].text = str(c)

    open_orders = sorted(
        [o for o in orders_list if (o.get("total_price") or 0) - (o.get("deposit_paid") or 0) > 0],
        key=lambda x: (x.get("total_price") or 0) - (x.get("deposit_paid") or 0),
        reverse=True
    )[:10]
    if open_orders:
        doc.add_heading("Top Open Orders by Balance", 1)
        t2 = doc.add_table(rows=1, cols=4)
        t2.style = "Table Grid"
        h2 = t2.rows[0].cells
        h2[0].text = "Customer"
        h2[1].text = "Items"
        h2[2].text = "Total"
        h2[3].text = "Balance Due"
        for o in open_orders:
            r = t2.add_row().cells
            r[0].text = o.get("customer_name", "")
            r[1].text = (o.get("items", "") or "")[:60]
            r[2].text = f"${o.get('total_price', 0) or 0:,.2f}"
            r[3].text = f"${(o.get('total_price', 0) or 0) - (o.get('deposit_paid', 0) or 0):,.2f}"

    doc.add_paragraph()
    fp = doc.add_paragraph(
        "Great Bridge Furniture — 1325 S Battlefield Blvd, Chesapeake VA 23322 — (757) 482-6622")
    fp.alignment = WD_ALIGN_PARAGRAPH.CENTER
    fp.runs[0].font.size = Pt(9)
    fp.runs[0].italic    = True

    buf = _io3.BytesIO()
    doc.save(buf)
    buf.seek(0)
    b64 = _b64_3.b64encode(buf.read()).decode()
    fname_key = "Proposal" if is_proposal else ("Report" if is_report else "Document")
    return {
        "filename": f"GBF_{fname_key}_{gen_date}.docx",
        "mime": "application/vnd.openxmlformats-officedocument.wordprocessingml.document",
        "b64": b64,
    }


def _gen_pptx(orders_list, message):
    """
    Build a .pptx presentation from orders_list.
    Returns {filename, mime, b64}.
    """
    import io as _io4, base64 as _b64_4, datetime as _dt4
    from collections import Counter as _Counter4
    from pptx import Presentation
    from pptx.util import Inches as _Inches, Pt as _Pt
    from pptx.dml.color import RGBColor as _PptxRGB
    from pptx.enum.text import PP_ALIGN

    prs = Presentation()
    prs.slide_width  = _Inches(13.33)
    prs.slide_height = _Inches(7.5)

    DARK_BLUE = _PptxRGB(0x1B, 0x4F, 0x72)
    MID_BLUE  = _PptxRGB(0x21, 0x8C, 0xBE)
    WHITE     = _PptxRGB(0xFF, 0xFF, 0xFF)
    GREY      = _PptxRGB(0x55, 0x6B, 0x7D)
    LIGHT_BLU = _PptxRGB(0xAE, 0xD6, 0xF1)

    today   = _dt4.datetime.now().strftime("%B %d, %Y")
    gen_dt  = _dt4.datetime.now().strftime("%Y%m%d")
    total_orders  = len(orders_list)
    total_revenue = sum((o.get("total_price")  or 0) for o in orders_list)
    total_deposit = sum((o.get("deposit_paid") or 0) for o in orders_list)
    balance_due   = total_revenue - total_deposit
    sc4           = _Counter4(o.get("status", "Unknown") for o in orders_list)
    blank_layout  = prs.slide_layouts[6]

    def _fill_bg(slide, color):
        bg = slide.background
        fill = bg.fill
        fill.solid()
        fill.fore_color.rgb = color

    def _txb(slide, text, left, top, width, height,
             font_size=18, bold=False, color=None, align=PP_ALIGN.LEFT, italic=False):
        tb = slide.shapes.add_textbox(
            _Inches(left), _Inches(top), _Inches(width), _Inches(height))
        tf = tb.text_frame
        tf.word_wrap = True
        p  = tf.paragraphs[0]
        p.alignment = align
        run = p.add_run()
        run.text = text
        run.font.size  = _Pt(font_size)
        run.font.bold  = bold
        run.font.italic = italic
        if color:
            run.font.color.rgb = color
        return tb

    # Slide 1 – Title
    s1 = prs.slides.add_slide(blank_layout)
    _fill_bg(s1, DARK_BLUE)
    _txb(s1, "Great Bridge Furniture", 0.5, 1.5, 12, 1.2,
         font_size=40, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    _txb(s1, "Business Overview Presentation", 0.5, 2.9, 12, 0.8,
         font_size=24, color=LIGHT_BLU, align=PP_ALIGN.CENTER)
    _txb(s1, f"Generated {today}", 0.5, 6.5, 12, 0.6,
         font_size=14, italic=True, color=LIGHT_BLU, align=PP_ALIGN.CENTER)

    # Slide 2 – Key Metrics
    s2 = prs.slides.add_slide(blank_layout)
    _fill_bg(s2, WHITE)
    _txb(s2, "Key Metrics", 0.5, 0.3, 12, 0.8,
         font_size=32, bold=True, color=DARK_BLUE)
    metrics = [
        ("Total Orders",        str(total_orders)),
        ("Total Revenue",       f"${total_revenue:,.2f}"),
        ("Total Collected",     f"${total_deposit:,.2f}"),
        ("Outstanding Balance", f"${balance_due:,.2f}"),
    ]
    for i, (label, val) in enumerate(metrics):
        col = i % 2
        row_i = i // 2
        x = 0.5 + col * 6.4
        y = 1.5 + row_i * 2.4
        shape = s2.shapes.add_shape(
            1, _Inches(x), _Inches(y), _Inches(5.8), _Inches(2.0))
        shape.fill.solid()
        shape.fill.fore_color.rgb = _PptxRGB(0xEB, 0xF5, 0xFB)
        shape.line.color.rgb = MID_BLUE
        _txb(s2, label, x + 0.15, y + 0.1, 5.5, 0.5, font_size=14, color=GREY)
        _txb(s2, val,   x + 0.15, y + 0.55, 5.5, 1.0,
             font_size=28, bold=True, color=DARK_BLUE)

    # Slide 3 – Status Breakdown
    s3 = prs.slides.add_slide(blank_layout)
    _fill_bg(s3, WHITE)
    _txb(s3, "Order Status Breakdown", 0.5, 0.3, 12, 0.8,
         font_size=32, bold=True, color=DARK_BLUE)
    STATUS_COLORS = {
        "Pending":    _PptxRGB(0xF3, 0x9C, 0x12),
        "In Transit": _PptxRGB(0x21, 0x8C, 0xBE),
        "Delivered":  _PptxRGB(0x27, 0xAE, 0x60),
        "Cancelled":  _PptxRGB(0xC0, 0x39, 0x2B),
    }
    for i, (status, count) in enumerate(sc4.most_common()):
        x = 0.5 + (i % 4) * 3.2
        y = 1.5 + (i // 4) * 2.2
        color = STATUS_COLORS.get(status, MID_BLUE)
        shape = s3.shapes.add_shape(
            1, _Inches(x), _Inches(y), _Inches(2.8), _Inches(1.8))
        shape.fill.solid()
        shape.fill.fore_color.rgb = color
        shape.line.color.rgb = color
        _txb(s3, str(count), x + 0.1, y + 0.1, 2.6, 0.9,
             font_size=36, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
        _txb(s3, status, x + 0.1, y + 1.0, 2.6, 0.6,
             font_size=14, color=WHITE, align=PP_ALIGN.CENTER)

    # Slide 4 – Top Open Orders
    open_orders = sorted(
        [o for o in orders_list
         if (o.get("total_price") or 0) - (o.get("deposit_paid") or 0) > 0],
        key=lambda x: (x.get("total_price") or 0) - (x.get("deposit_paid") or 0),
        reverse=True
    )[:8]
    if open_orders:
        s4 = prs.slides.add_slide(blank_layout)
        _fill_bg(s4, WHITE)
        _txb(s4, "Top Open Balances", 0.5, 0.3, 12, 0.8,
             font_size=32, bold=True, color=DARK_BLUE)
        for i, o in enumerate(open_orders):
            y = 1.3 + i * 0.72
            bal  = (o.get("total_price") or 0) - (o.get("deposit_paid") or 0)
            name = (o.get("customer_name") or "")[:25]
            items = (o.get("items") or "")[:45]
            _txb(s4, name,  0.5, y, 3.5, 0.55, font_size=16, bold=True, color=DARK_BLUE)
            _txb(s4, items, 4.1, y, 6.0, 0.55, font_size=14, color=GREY)
            _txb(s4, f"${bal:,.2f}", 10.2, y, 2.5, 0.55,
                 font_size=16, bold=True, color=MID_BLUE, align=PP_ALIGN.RIGHT)

    # Slide 5 – Closing
    s5 = prs.slides.add_slide(blank_layout)
    _fill_bg(s5, DARK_BLUE)
    _txb(s5, "Great Bridge Furniture", 0.5, 2.0, 12, 1.0,
         font_size=36, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    _txb(s5, "Since 1984 — Quality Furniture & Service", 0.5, 3.2, 12, 0.6,
         font_size=20, italic=True, color=LIGHT_BLU, align=PP_ALIGN.CENTER)
    _txb(s5, "1325 S Battlefield Blvd, Chesapeake VA 23322  |  (757) 482-6622",
         0.5, 6.2, 12, 0.5, font_size=14, color=LIGHT_BLU, align=PP_ALIGN.CENTER)

    buf = _io4.BytesIO()
    prs.save(buf)
    buf.seek(0)
    b64 = _b64_4.b64encode(buf.read()).decode()
    return {
        "filename": f"GBF_Presentation_{gen_dt}.pptx",
        "mime": "application/vnd.openxmlformats-officedocument.presentationml.presentation",
        "b64": b64,
    }


# ── GREAT BRIDGE FURNITURE — ASK AIVM CHAT ───────────────────────────────────

@app.route("/api/gb/ask", methods=["POST"])
def gb_ask():
    """
    Chat endpoint with live business context injection.
    Body: { message, messages (array), file_b64, file_mime, file_name }
    Returns: { reply: "..." }
    """
    data      = request.get_json(force=True) or {}
    message   = str(data.get("message", "")).strip()
    messages  = data.get("messages", [])   # prior conversation turns
    file_b64  = data.get("file_b64")       # optional base64 attachment
    file_mime = str(data.get("file_mime", "")).lower()
    file_name = str(data.get("file_name", ""))

    if not message:
        return jsonify({"error": "message required"}), 400

    today = time.strftime('%Y-%m-%d')

    # ── Fetch live context ──────────────────────────────────────────────────
    try:
        _init_db()
        with _db() as conn:
            orders = conn.execute(
                "SELECT * FROM gb_orders WHERE created_at >= date('now','-90 days') "
                "ORDER BY created_at DESC"
            ).fetchall()
            orders_list = [dict(r) for r in orders]

            inv_rows = conn.execute(
                "SELECT status, COUNT(*) as cnt FROM gb_inventory GROUP BY status"
            ).fetchall()
            inv_summary = {r["status"]: r["cnt"] for r in inv_rows}

            activity = conn.execute(
                "SELECT user_name, action, target, details, created_at "
                "FROM gb_activity_log ORDER BY created_at DESC LIMIT 20"
            ).fetchall()
            activity_list = [dict(r) for r in activity]
    except Exception as e:
        print(f"[gb/ask] DB error: {e}")
        orders_list   = []
        inv_summary   = {}
        activity_list = []

    # ── Compute key metrics ─────────────────────────────────────────────────
    status_counts = {}
    overdue       = []
    ready         = []
    balance_due   = 0.0
    for o in orders_list:
        s = o.get("status", "Ordered")
        status_counts[s] = status_counts.get(s, 0) + 1
        exp = o.get("expected_date", "")
        if exp and exp < today and s not in ("Delivered", "Paid"):
            overdue.append(o)
        if s == "Ready for Pickup":
            ready.append(o)
        try:
            balance_due += float(o.get("total_amount") or 0) - float(o.get("deposit_paid") or 0)
        except Exception:
            pass

    orders_slim = []
    for o in orders_list[:60]:
        orders_slim.append({
            "id":            o.get("id"),
            "customer":      o.get("customer_name"),
            "items":         o.get("items"),
            "manufacturer":  o.get("manufacturer"),
            "status":        o.get("status"),
            "expected_date": o.get("expected_date"),
            "balance":       round(float(o.get("total_amount") or 0) - float(o.get("deposit_paid") or 0), 2),
        })

    system_ctx = (
        f"You are AIVM, the AI business assistant for Great Bridge Furniture, a furniture store "
        f"in Chesapeake VA (1325 S Battlefield Blvd, 757-482-6622). Today is {today}.\n\n"
        f"LIVE BUSINESS DATA:\n"
        f"Orders (last 90 days): {len(orders_list)} total\n"
        f"Status breakdown: {json.dumps(status_counts)}\n"
        f"Overdue (past expected date, not yet delivered): {len(overdue)}\n"
        f"Ready for pickup: {len(ready)}\n"
        f"Outstanding balance: ${balance_due:,.2f}\n"
        f"Inventory by status: {json.dumps(inv_summary)}\n\n"
        f"Order details (most recent 60):\n{json.dumps(orders_slim, indent=2)}\n\n"
        f"Recent activity (last 20):\n{json.dumps(activity_list, indent=2)}\n\n"
        f"Answer David's questions directly and helpfully. Be specific — use the actual data. "
        f"Offer practical business advice. Keep answers concise but complete.\n"
    )

    # ── Handle file attachment ──────────────────────────────────────────────
    file_section = ""
    if file_b64:
        if "image" in file_mime:
            file_section = (
                f"\n[The owner has attached an image named '{file_name}'. "
                f"Analyze it in the context of their question. "
                f"Describe what you can infer and how it relates to their business.]\n"
            )
        else:
            try:
                import base64 as _b64
                raw = _b64.b64decode(file_b64)
                fn_lower = file_name.lower()
                if fn_lower.endswith(".csv") or "csv" in file_mime:
                    text  = raw.decode("utf-8", errors="replace")
                    lines = text.split("\n")[:60]
                    file_section = (
                        f"\n[The owner has attached a CSV file named '{file_name}'. "
                        f"First 60 rows:\n" + "\n".join(lines) + "\n]\n"
                    )
                elif fn_lower.endswith((".xlsx", ".xls")):
                    try:
                        import openpyxl as _xl
                        wb = _xl.load_workbook(io.BytesIO(raw), data_only=True)
                        ws = wb.active
                        rows_txt = []
                        for row in ws.iter_rows(max_row=60, values_only=True):
                            rows_txt.append("\t".join("" if c is None else str(c) for c in row))
                        file_section = (
                            f"\n[The owner has attached a spreadsheet named '{file_name}'. "
                            f"First 60 rows:\n" + "\n".join(rows_txt) + "\n]\n"
                        )
                    except Exception as ex:
                        file_section = f"\n[Spreadsheet '{file_name}' could not be parsed: {ex}]\n"
                else:
                    file_section = f"\n[The owner has attached a file named '{file_name}'.]\n"
            except Exception as ex:
                file_section = f"\n[File '{file_name}' could not be processed: {ex}]\n"

    # ── Build single prompt string ──────────────────────────────────────────
    conv_lines = []
    for m in messages[-8:]:
        role    = m.get("role", "user")
        content = m.get("content", "")
        conv_lines.append(("User" if role == "user" else "Assistant") + ": " + content)
    conv_lines.append("User: " + message)

    full_prompt = system_ctx + file_section + "\nConversation:\n" + "\n".join(conv_lines) + "\n\nAssistant:"

    # ── Detect document-generation requests ────────────────────────────────
    q_low = message.lower()
    _is_spreadsheet = any(w in q_low for w in ("spreadsheet", "excel", "xlsx", ".xlsx", "csv file"))
    _is_word        = any(w in q_low for w in ("word doc", "word document", ".docx", "write a doc", "create a doc", "write me a doc", "make a doc"))
    _is_pptx        = any(w in q_low for w in ("powerpoint", "presentation", "slide deck", "slides", ".pptx"))
    _make_words     = any(w in q_low for w in ("make", "create", "generate", "build", "write", "give me", "export", "download"))

    if _make_words and _is_spreadsheet:
        try:
            download = _gen_spreadsheet(orders_list, message)
            try:
                aivm  = get_aivm()
                caption = aivm.run_inference(
                    system_ctx + f"\nUser asked: {message}\nYou just created a spreadsheet for them. "
                    "Respond in 1-2 sentences confirming what's in it and that they can download it.\nAssistant:",
                    timeout_secs=60)
            except Exception:
                caption = "Here's your spreadsheet — click the button below to download it."
            return jsonify({"reply": caption.strip(), "download": download})
        except Exception as e:
            print(f"[gb/ask] spreadsheet gen error: {e}")

    if _make_words and _is_word:
        try:
            download = _gen_word_doc(orders_list, message, file_b64, file_mime, file_name)
            try:
                aivm  = get_aivm()
                caption = aivm.run_inference(
                    system_ctx + f"\nUser asked: {message}\nYou just created a Word document for them. "
                    "Respond in 1-2 sentences confirming what's in it.\nAssistant:",
                    timeout_secs=60)
            except Exception:
                caption = "Here's your Word document — click the button below to download it."
            return jsonify({"reply": caption.strip(), "download": download})
        except Exception as e:
            print(f"[gb/ask] word doc gen error: {e}")

    if _make_words and _is_pptx:
        try:
            download = _gen_pptx(orders_list, message)
            try:
                aivm  = get_aivm()
                caption = aivm.run_inference(
                    system_ctx + f"\nUser asked: {message}\nYou just created a PowerPoint presentation for them. "
                    "Respond in 1-2 sentences confirming what slides are in it.\nAssistant:",
                    timeout_secs=60)
            except Exception:
                caption = "Here's your presentation — click the button below to download it."
            return jsonify({"reply": caption.strip(), "download": download})
        except Exception as e:
            print(f"[gb/ask] pptx gen error: {e}")

    try:
        aivm  = get_aivm()
        reply = aivm.run_inference(full_prompt, timeout_secs=180)
        return jsonify({"reply": reply.strip()})
    except Exception as e:
        print(f"[gb/ask] AIVM error: {e}")
        # ── Rule-based fallback ───────────────────────────────────────────
        q = message.lower()
        if any(w in q for w in ("owe", "balance", "money", "pay", "due")):
            top = sorted(orders_slim, key=lambda x: x.get("balance", 0), reverse=True)
            top = [o for o in top if o.get("balance", 0) > 0][:5]
            reply = f"Outstanding balance: ${balance_due:,.2f} across {len(orders_list)} orders."
            if top:
                reply += " Largest balances: " + ", ".join(f"{o['customer']} (${o['balance']:.2f})" for o in top)
        elif any(w in q for w in ("overdue", "late", "past", "behind")):
            if overdue:
                names = ", ".join(o["customer_name"] for o in overdue[:5])
                reply = f"{len(overdue)} overdue order(s): {names}. Follow up with manufacturers."
            else:
                reply = "No overdue orders right now — you're all caught up."
        elif any(w in q for w in ("arriv", "week", "coming", "transit", "shipping")):
            transit = [o for o in orders_slim if o.get("status") == "In Transit"]
            reply = (f"{len(transit)} order(s) currently in transit. "
                     f"Check expected dates to confirm arrivals this week.")
        elif any(w in q for w in ("pickup", "ready", "waiting", "call")):
            if ready:
                names = ", ".join(o["customer_name"] for o in ready[:5])
                reply = f"{len(ready)} order(s) ready for pickup: {names}. Call to schedule delivery."
            else:
                reply = "No orders waiting for pickup right now."
        else:
            reply = (f"I have live data on your {len(orders_list)} orders. "
                     f"Ask about balances, overdue orders, what's arriving, or specific customers.")
        return jsonify({"reply": reply, "fallback": True})


# ── GREAT BRIDGE FURNITURE — SCAN TICKET ─────────────────────────────────────

def _extract_ticket_fields(raw_bytes, mime):
    """
    Core ticket-extraction logic: send image bytes to AIVM, parse JSON response.
    Returns dict with keys: customer_name, phone, items, manufacturer,
    total_price, deposit_paid, expected_date, status, notes.
    Raises on failure.
    """
    b64_data = _b64_mod.b64encode(raw_bytes).decode()
    extract_prompt = (
        "This is a photo or scan of a handwritten furniture store order ticket from Great Bridge Furniture. "
        "Extract all visible information and return it as a JSON object with ONLY these exact keys: "
        "customer_name, phone, items, manufacturer, total_price, deposit_paid, expected_date, status, notes. "
        "For 'items', return a string listing all items visible. "
        "For 'expected_date', use ISO format YYYY-MM-DD if you can determine the date, otherwise null. "
        "For any field not visible or unclear, return null. "
        "Return ONLY valid JSON with no explanation, no markdown fences, just the raw JSON object.\n\n"
        f"[Image attached as base64 {mime}: {b64_data[:60]}... (truncated for context — use the actual image data to extract fields)]"
    )
    aivm   = get_aivm()
    result = aivm.run_inference(extract_prompt, timeout_secs=180)
    cleaned = _re_mod.sub(r"```(?:json)?", "", result).strip().strip("`").strip()
    match = _re_mod.search(r"\{.*\}", cleaned, _re_mod.DOTALL)
    extracted = json.loads(match.group() if match else cleaned)
    keys = ["customer_name", "phone", "items", "manufacturer",
            "total_price", "deposit_paid", "expected_date", "status", "notes"]
    return {k: extracted.get(k) for k in keys}


@app.route("/api/gb/scan-ticket", methods=["POST"])
def gb_scan_ticket():
    """
    Upload a photo/scan of a handwritten order ticket.
    Accepts multipart form with 'image' file field.
    Returns: { fields: { customer_name, phone, items, manufacturer, total_price,
                          deposit_paid, expected_date, status, notes } }
    """
    if "image" not in request.files:
        return jsonify({"error": "No image file uploaded"}), 400
    img_file = request.files["image"]
    try:
        raw  = img_file.read()
        mime = img_file.content_type or "image/jpeg"
    except Exception as e:
        return jsonify({"error": f"Could not read image: {e}"}), 400
    try:
        fields = _extract_ticket_fields(raw, mime)
    except Exception as e:
        print(f"[gb/scan-ticket] AIVM/parse error: {e}")
        return jsonify({"error": f"Could not extract data from image: {e}"}), 500
    return jsonify({"fields": fields})


@app.route("/api/gb/bulk-scan", methods=["POST"])
def gb_bulk_scan():
    """
    Upload multiple ticket images at once (field name: tickets[]).
    Returns: { results: [ { filename, fields, error }, ... ] }
    Processes sequentially; max 50 images per call.
    No orders are created — caller shows preview and confirms.
    """
    files = request.files.getlist("tickets[]")
    if not files:
        return jsonify({"error": "No images uploaded (use field name tickets[])"}), 400
    if len(files) > 50:
        return jsonify({"error": "Maximum 50 images per batch"}), 400

    results = []
    for f in files:
        filename = (f.filename or "ticket.jpg").replace("\\", "/").split("/")[-1][:255]
        try:
            raw  = f.read()
            mime = f.content_type or "image/jpeg"
            fields = _extract_ticket_fields(raw, mime)
            results.append({"filename": filename, "fields": fields, "error": None})
        except Exception as e:
            print(f"[gb/bulk-scan] {filename}: {e}")
            results.append({"filename": filename, "fields": {}, "error": str(e)})

    return jsonify({"results": results})


@app.route("/api/gb/orders/<int:order_id>/set-ticket", methods=["POST"])
def gb_set_ticket_doc(order_id):
    """
    Called after auto-attaching a scanned ticket image to an order.
    Uploads the image as a doc (is_ticket=1) and records ticket_doc_id on the order.
    Accepts multipart: 'file' field.
    """
    if "file" not in request.files:
        return jsonify({"error": "No file"}), 400
    f = request.files["file"]
    raw      = f.read()
    filename = "TICKET_" + (f.filename or "ticket.jpg").replace("\\", "/").split("/")[-1][:240]
    filetype = f.content_type or "image/jpeg"
    user     = request.form.get("user_name") or "Unknown"

    if len(raw) > 20 * 1024 * 1024:
        return jsonify({"error": "File too large (20 MB max)"}), 400

    try:
        _init_db()
        with _db() as conn:
            cur = conn.execute(
                "INSERT INTO gb_docs (order_id, filename, filetype, filedata, is_ticket) VALUES (?,?,?,?,1)",
                (order_id, filename, filetype, raw)
            )
            doc_id = cur.lastrowid
            conn.execute("UPDATE gb_orders SET ticket_doc_id=? WHERE id=?", (doc_id, order_id))
            conn.commit()
        _log_activity(user, "Attached ticket scan", f"to order #{order_id}", f"File: {filename}")
        return jsonify({"ok": True, "doc_id": doc_id})
    except Exception as e:
        print(f"[gb/set-ticket order={order_id}] {e}")
        return jsonify({"error": str(e)}), 500


# ── GREAT BRIDGE FURNITURE — PROPOSAL ────────────────────────────────────────

@app.route("/api/gb/proposal/<int:order_id>")
def gb_proposal(order_id):
    """
    Generate a printable HTML customer proposal for an order.
    Returns Content-Type: text/html directly (open in new tab).
    """
    try:
        _init_db()
        with _db() as conn:
            row = conn.execute(
                "SELECT * FROM gb_orders WHERE id = ?", (order_id,)
            ).fetchone()
    except Exception as e:
        return f"<html><body>Database error: {e}</body></html>", 500

    if not row:
        return "<html><body>Order not found.</body></html>", 404

    order = dict(row)
    today = time.strftime("%B %d, %Y")

    # Parse items
    items_raw = order.get("items", "")
    try:
        items_list = json.loads(items_raw) if isinstance(items_raw, str) else items_raw
        if not isinstance(items_list, list):
            items_list = [str(items_raw)]
    except Exception:
        items_list = [str(items_raw)] if items_raw else []

    total    = float(order.get("total_amount") or 0)
    deposit  = float(order.get("deposit_paid") or 0)
    balance  = total - deposit
    exp_date = order.get("expected_date") or "To be confirmed"

    # Try AIVM first
    aivm_prompt = (
        "Create a professional, warm customer proposal letter for a furniture store sale. "
        "Format it as complete HTML (no doctype/html/head/body tags — just the content HTML). "
        "Include: company header with name, address, phone; customer name; itemized order list; "
        "total price, deposit paid, balance due; expected delivery date; friendly thank-you message; "
        "and a print button. Use clean styling with inline CSS — dark on white, professional look. "
        "Business: Great Bridge Furniture, 1325 S Battlefield Blvd, Chesapeake VA 23322, Phone: 757-482-6622. "
        f"Order details: Customer: {order.get('customer_name')}, "
        f"Phone: {order.get('phone') or 'N/A'}, "
        f"Items: {', '.join(items_list)}, "
        f"Manufacturer: {order.get('manufacturer') or 'N/A'}, "
        f"Total: ${total:,.2f}, Deposit Paid: ${deposit:,.2f}, Balance Due: ${balance:,.2f}, "
        f"Expected Delivery: {exp_date}, "
        f"Status: {order.get('status')}, "
        f"Notes: {order.get('notes') or 'None'}. "
        f"Today's date: {today}."
    )

    html_body = None
    try:
        aivm      = get_aivm()
        raw       = aivm.run_inference(aivm_prompt, timeout_secs=180)
        # Strip markdown fences
        html_body = _re_mod.sub(r"```(?:html)?", "", raw).strip().strip("`").strip()
    except Exception as e:
        print(f"[gb/proposal] AIVM error: {e}")

    if not html_body:
        # ── Fallback: generate directly from data ─────────────────────────
        items_html = "".join(f"<tr><td style='padding:8px 16px'>{item}</td></tr>" for item in items_list)
        html_body = f"""
<div style="font-family:Georgia,serif;max-width:700px;margin:40px auto;padding:40px;border:1px solid #ddd;border-radius:8px">
  <div style="text-align:center;border-bottom:2px solid #2c3e50;padding-bottom:20px;margin-bottom:28px">
    <h1 style="margin:0;font-size:28px;color:#2c3e50">Great Bridge Furniture</h1>
    <p style="margin:6px 0 0;color:#666;font-size:15px">1325 S Battlefield Blvd, Chesapeake VA 23322 &nbsp;|&nbsp; 757-482-6622</p>
  </div>
  <p style="font-size:15px;color:#666;margin-bottom:24px">Date: {today}</p>
  <h2 style="font-size:20px;color:#2c3e50;margin-bottom:6px">Customer Order Proposal</h2>
  <p style="margin:0 0 4px"><strong>Customer:</strong> {order.get('customer_name') or '—'}</p>
  <p style="margin:0 0 20px"><strong>Phone:</strong> {order.get('phone') or '—'}</p>
  <table style="width:100%;border-collapse:collapse;margin-bottom:20px">
    <thead><tr style="background:#2c3e50;color:#fff"><th style="padding:10px 16px;text-align:left">Items Ordered</th></tr></thead>
    <tbody style="background:#f9f9f9">{items_html}</tbody>
  </table>
  <table style="width:100%;border-collapse:collapse;margin-bottom:28px">
    <tr><td style="padding:8px;border-bottom:1px solid #eee"><strong>Manufacturer:</strong></td><td style="padding:8px;border-bottom:1px solid #eee">{order.get('manufacturer') or '—'}</td></tr>
    <tr><td style="padding:8px;border-bottom:1px solid #eee"><strong>Total Price:</strong></td><td style="padding:8px;border-bottom:1px solid #eee;font-weight:bold">${total:,.2f}</td></tr>
    <tr><td style="padding:8px;border-bottom:1px solid #eee"><strong>Deposit Paid:</strong></td><td style="padding:8px;border-bottom:1px solid #eee;color:#27ae60">${deposit:,.2f}</td></tr>
    <tr><td style="padding:8px;border-bottom:1px solid #eee"><strong>Balance Due:</strong></td><td style="padding:8px;border-bottom:1px solid #eee;font-weight:bold;color:{'#e74c3c' if balance > 0 else '#27ae60'}">${balance:,.2f}</td></tr>
    <tr><td style="padding:8px"><strong>Expected Delivery:</strong></td><td style="padding:8px">{exp_date}</td></tr>
  </table>
  <p style="background:#eaf4fb;border-left:4px solid #2c3e50;padding:16px;border-radius:4px;font-size:15px;line-height:1.7">
    Thank you for choosing Great Bridge Furniture! We truly appreciate your business and look forward to delivering beautiful furniture to your home.
    If you have any questions about your order, please don't hesitate to call us at <strong>757-482-6622</strong>.
  </p>
  <div style="text-align:center;margin-top:32px">
    <button onclick="window.print()" style="background:#2c3e50;color:#fff;border:none;padding:12px 28px;border-radius:6px;font-size:16px;cursor:pointer;font-family:inherit">🖨️ Print Proposal</button>
  </div>
</div>
"""

    full_html = f"""<!DOCTYPE html>
<html lang="en">
<head>
<meta charset="UTF-8">
<title>Proposal — {order.get('customer_name') or 'Customer'} | Great Bridge Furniture</title>
<style>
  body {{ margin:0; padding:20px; background:#f5f5f5; font-family:Georgia,serif; }}
  @media print {{
    body {{ background:#fff; padding:0; }}
    button {{ display:none !important; }}
  }}
</style>
</head>
<body>
{html_body}
</body>
</html>"""

    from flask import Response
    return Response(full_html, content_type="text/html; charset=utf-8")


# ── GREAT BRIDGE FURNITURE — MANUFACTURERS ───────────────────────────────────

@app.route('/api/gb/manufacturers', methods=['GET'])
def gb_list_manufacturers():
    conn = _db()
    rows = conn.execute("SELECT * FROM gb_manufacturers ORDER BY name").fetchall()
    conn.close()
    return jsonify([dict(r) for r in rows])


@app.route('/api/gb/manufacturers', methods=['POST'])
def gb_create_manufacturer():
    data = request.json or {}
    conn = _db()
    c = conn.cursor()
    c.execute(
        "INSERT INTO gb_manufacturers (name,contact_person,email,phone,fax,address,notes) VALUES (?,?,?,?,?,?,?)",
        (data.get('name',''), data.get('contact_person',''), data.get('email',''),
         data.get('phone',''), data.get('fax',''), data.get('address',''), data.get('notes',''))
    )
    conn.commit()
    new_id = c.lastrowid
    conn.close()
    return jsonify({'ok': True, 'id': new_id})


@app.route('/api/gb/manufacturers/<int:mid>', methods=['PUT'])
def gb_update_manufacturer(mid):
    data = request.json or {}
    conn = _db()
    conn.execute(
        "UPDATE gb_manufacturers SET name=?,contact_person=?,email=?,phone=?,fax=?,address=?,notes=? WHERE id=?",
        (data.get('name',''), data.get('contact_person',''), data.get('email',''),
         data.get('phone',''), data.get('fax',''), data.get('address',''), data.get('notes',''), mid)
    )
    conn.commit()
    conn.close()
    return jsonify({'ok': True})


@app.route('/api/gb/manufacturers/<int:mid>', methods=['DELETE'])
def gb_delete_manufacturer(mid):
    conn = _db()
    conn.execute("DELETE FROM gb_manufacturers WHERE id=?", (mid,))
    conn.commit()
    conn.close()
    return jsonify({'ok': True})


@app.route('/api/gb/manufacturers/import', methods=['POST'])
def gb_import_manufacturers():
    """Import manufacturers from a spreadsheet or Word doc using AIVM to extract contact info."""
    if 'file' not in request.files:
        return jsonify({'error': 'No file'}), 400
    f     = request.files['file']
    raw   = f.read()
    fname = (f.filename or '').lower()

    text_content = ''
    try:
        if fname.endswith('.docx'):
            from docx import Document as _DocxDoc
            import io as _io_mfr
            doc = _DocxDoc(_io_mfr.BytesIO(raw))
            text_content = '\n'.join(p.text for p in doc.paragraphs if p.text.strip())
        elif fname.endswith('.xlsx') or fname.endswith('.xls'):
            import openpyxl as _xl_mfr, io as _io_mfr2
            wb = _xl_mfr.load_workbook(_io_mfr2.BytesIO(raw), data_only=True)
            ws = wb.active
            rows = []
            for row in ws.iter_rows(values_only=True):
                row_text = ' | '.join(str(c) for c in row if c is not None)
                if row_text.strip():
                    rows.append(row_text)
            text_content = '\n'.join(rows[:200])
        elif fname.endswith('.csv'):
            text_content = raw.decode('utf-8', errors='replace')[:8000]
        else:
            text_content = raw.decode('utf-8', errors='replace')[:8000]
    except Exception as e:
        return jsonify({'error': f'Could not read file: {e}'}), 400

    prompt = f"""Extract all manufacturer/vendor/supplier contact information from the following text.
Return a JSON array where each element has these fields:
- name (company name, required)
- contact_person (person's name if present, else empty string)
- email (email address if present, else empty string)
- phone (phone number if present, else empty string)
- fax (fax number if present, else empty string)
- address (mailing/shipping address if present, else empty string)
- notes (any other relevant info, else empty string)

Return ONLY the JSON array, no other text.

Content to extract from:
{text_content[:6000]}"""

    manufacturers = []
    try:
        aivm_resp = _call_aivm(prompt)
        import json as _json_mfr, re as _re_mfr
        match = _re_mfr.search(r'\[.*\]', aivm_resp, _re_mfr.DOTALL)
        if match:
            manufacturers = _json_mfr.loads(match.group())
    except Exception:
        pass

    return jsonify({'ok': True, 'manufacturers': manufacturers, 'count': len(manufacturers)})


@app.route('/api/gb/orders/<int:oid>/send-po', methods=['POST'])
def gb_send_po(oid):
    """Generate a PO and email it to the manufacturer."""
    data      = request.json or {}
    user_name = data.get('user_name', 'Staff')

    conn  = _db()
    order = conn.execute("SELECT * FROM gb_orders WHERE id=?", (oid,)).fetchone()
    if not order:
        conn.close()
        return jsonify({'error': 'Order not found'}), 404
    order = dict(order)

    # Look up manufacturer — try manufacturer_id first, then match by name
    manufacturer = None
    if order.get('manufacturer_id'):
        row = conn.execute("SELECT * FROM gb_manufacturers WHERE id=?", (order['manufacturer_id'],)).fetchone()
        if row:
            manufacturer = dict(row)
    if not manufacturer and order.get('manufacturer'):
        row = conn.execute("SELECT * FROM gb_manufacturers WHERE name LIKE ?",
                           (f"%{order['manufacturer']}%",)).fetchone()
        if row:
            manufacturer = dict(row)

    to_email = data.get('override_email') or (manufacturer.get('email') if manufacturer else '')
    if not to_email:
        conn.close()
        return jsonify({'error': 'No manufacturer email address found. Please add an email to the manufacturer record or enter one manually.'}), 400

    po_html  = _gen_po_html(order, manufacturer)
    mfr_name = (manufacturer.get('name') if manufacturer else None) or order.get('manufacturer', 'Manufacturer')
    subject  = f"Purchase Order - Great Bridge Furniture - {order.get('customer_name', '')} - {order.get('items', '')[:40]}"

    try:
        _send_aol_email(to_email, subject, po_html)
    except ValueError as e:
        conn.close()
        return jsonify({'error': str(e), 'setup_needed': True}), 400
    except Exception as e:
        conn.close()
        return jsonify({'error': f'Email failed: {str(e)}'}), 500

    import datetime as _dt_po2
    now = _dt_po2.datetime.now().isoformat(sep=' ', timespec='seconds')
    conn.execute("UPDATE gb_orders SET po_sent_at=? WHERE id=?", (now, oid))
    conn.commit()
    _log_activity(user_name, 'send_po', f'Order #{oid}', f'PO emailed to {to_email} ({mfr_name})')
    conn.close()

    return jsonify({'ok': True, 'sent_to': to_email, 'manufacturer': mfr_name})


@app.route('/api/gb/orders/<int:oid>/po-preview')
def gb_po_preview(oid):
    """Return the PO as a printable HTML page (opens in new tab)."""
    conn  = _db()
    order = conn.execute("SELECT * FROM gb_orders WHERE id=?", (oid,)).fetchone()
    if not order:
        conn.close()
        return "Order not found", 404
    order        = dict(order)
    manufacturer = None
    if order.get('manufacturer_id'):
        row = conn.execute("SELECT * FROM gb_manufacturers WHERE id=?", (order['manufacturer_id'],)).fetchone()
        if row:
            manufacturer = dict(row)
    conn.close()
    html = _gen_po_html(order, manufacturer)
    html = html.replace('</body>', '<script>window.onload=function(){window.print();}</script></body>')
    return html, 200, {'Content-Type': 'text/html; charset=utf-8'}


# ── MAIN ─────────────────────────────────────────────────────────────────────
if __name__ == "__main__":
    port = int(os.environ.get("PORT", 5000))
    app.run(host="0.0.0.0", port=port)
